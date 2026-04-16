from fastapi import FastAPI, APIRouter, UploadFile, File, HTTPException
from fastapi.responses import StreamingResponse
from dotenv import load_dotenv
from starlette.middleware.cors import CORSMiddleware
from motor.motor_asyncio import AsyncIOMotorClient
import os
import logging
from pathlib import Path
from pydantic import BaseModel, Field, ConfigDict
from typing import List, Optional, Union, Any
import uuid
from datetime import datetime, timezone
import httpx
from bs4 import BeautifulSoup
import PyPDF2
import io
import base64
from docx import Document as DocxDocument
from pptx import Presentation as PptxPresentation
import openpyxl
from emergentintegrations.llm.chat import LlmChat, UserMessage, ImageContent
from emergentintegrations.llm.openai.image_generation import OpenAIImageGeneration

ROOT_DIR = Path(__file__).parent
load_dotenv(ROOT_DIR / '.env')

# MongoDB connection
mongo_url = os.environ['MONGO_URL']
client = AsyncIOMotorClient(mongo_url)
db = client[os.environ['DB_NAME']]

# Emergent LLM Key
EMERGENT_LLM_KEY = os.environ.get('EMERGENT_LLM_KEY')

# Create the main app without a prefix
app = FastAPI()

# Create a router with the /api prefix
api_router = APIRouter(prefix="/api")

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# ─── MODELS ─────────────────────────────────────────────────────────────────

class StatusCheck(BaseModel):
    model_config = ConfigDict(extra="ignore")
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    client_name: str
    timestamp: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))

class StatusCheckCreate(BaseModel):
    client_name: str

class Source(BaseModel):
    model_config = ConfigDict(extra="ignore")
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    title: str
    type: str  # pdf, url, txt, etc.
    content: str  # extracted text content
    chunks: int = 0
    size: str = ""
    status: str = "processing"  # processing, indexed, error
    url: Optional[str] = None
    created_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    notebook_id: str = "default"

class SourceCreate(BaseModel):
    title: str
    type: str
    url: Optional[str] = None
    notebook_id: str = "default"

class GenerateRequest(BaseModel):
    output_type: str  # slides, report, mindmap, flashcards, quiz, audio, video, infographic, datatable
    notebook_id: str = "default"
    title: str = ""

class GenerateResponse(BaseModel):
    id: str
    type: str
    title: str
    content: str
    slides_data: Optional[Union[List[dict], dict]] = None  # For PPTX/canvas generation (list for slides/report/infographic/flashcards, dict for mindmap/quiz/datatable)
    theme: Optional[str] = None  # Theme for styling

# ─── HELPER FUNCTIONS ───────────────────────────────────────────────────────

async def extract_text_from_pdf(file_content: bytes) -> str:
    """Extract text from PDF file"""
    try:
        pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
        text = ""
        for page in pdf_reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
        return text.strip()
    except Exception as e:
        logger.error(f"PDF extraction error: {e}")
        return ""

async def extract_text_from_docx(file_content: bytes) -> str:
    """Extract text from DOCX file"""
    try:
        doc = DocxDocument(io.BytesIO(file_content))
        text_parts = []
        
        # Extract text from paragraphs
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text.strip())
        
        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    text_parts.append(" | ".join(row_text))
        
        return "\n".join(text_parts)
    except Exception as e:
        logger.error(f"DOCX extraction error: {e}")
        return ""

async def extract_text_from_pptx(file_content: bytes) -> str:
    """Extract text from PowerPoint PPTX file"""
    try:
        prs = PptxPresentation(io.BytesIO(file_content))
        text_parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_texts.append(text)
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if row_text:
                            slide_texts.append(" | ".join(row_text))
            if slide_texts:
                text_parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_texts))
        return "\n\n".join(text_parts)
    except Exception as e:
        logger.error(f"PPTX extraction error: {e}")
        return ""

async def extract_text_from_excel(file_content: bytes) -> str:
    """Extract text from Excel XLSX file"""
    try:
        wb = openpyxl.load_workbook(io.BytesIO(file_content), read_only=True, data_only=True)
        text_parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            sheet_rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(c.strip() for c in cells):
                    sheet_rows.append(" | ".join(c for c in cells if c.strip()))
            if sheet_rows:
                text_parts.append(f"--- Sheet: {sheet_name} ---\n" + "\n".join(sheet_rows))
        wb.close()
        return "\n\n".join(text_parts)
    except Exception as e:
        logger.error(f"Excel extraction error: {e}")
        return ""

async def extract_text_from_image(file_content: bytes, filename: str) -> str:
    """Extract text/description from image using AI vision"""
    try:
        image_b64 = base64.b64encode(file_content).decode('utf-8')

        chat = LlmChat(
            api_key=EMERGENT_LLM_KEY,
            session_id=str(uuid.uuid4()),
            system_message="You are an expert at extracting information from images. Extract ALL visible text exactly as it appears. If there is no text, provide a comprehensive description of the image contents, data, charts, diagrams, or visual information."
        ).with_model("openai", "gpt-5.1")

        img_content = ImageContent(image_base64=image_b64)
        user_message = UserMessage(
            text="Extract all text from this image. If it contains charts, tables, or diagrams, describe the data in detail. Be thorough.",
            file_contents=[img_content]
        )
        response = await chat.send_message(user_message)
        return response if response else ""
    except Exception as e:
        logger.error(f"Image extraction error: {e}")
        return ""

async def extract_text_from_url(url: str) -> str:
    """Extract text content from URL"""
    try:
        async with httpx.AsyncClient(timeout=30.0) as client:
            response = await client.get(url, follow_redirects=True)
            response.raise_for_status()
            soup = BeautifulSoup(response.text, 'html.parser')
            
            # Remove script and style elements
            for script in soup(["script", "style", "nav", "footer", "header"]):
                script.decompose()
            
            # Get text
            text = soup.get_text(separator='\n', strip=True)
            # Clean up whitespace
            lines = (line.strip() for line in text.splitlines())
            text = '\n'.join(line for line in lines if line)
            return text[:50000]  # Limit to 50k chars
    except Exception as e:
        logger.error(f"URL extraction error: {e}")
        return ""

def count_chunks(text: str, chunk_size: int = 500) -> int:
    """Count approximate chunks for indexing"""
    words = text.split()
    return max(1, len(words) // chunk_size)


def split_into_chunks(text: str, chunk_size: int = 400, overlap: int = 50) -> List[str]:
    """Split text into overlapping word-based chunks for retrieval"""
    words = text.split()
    if len(words) <= chunk_size:
        return [text] if text.strip() else []
    chunks = []
    start = 0
    while start < len(words):
        end = min(start + chunk_size, len(words))
        chunks.append(' '.join(words[start:end]))
        start += chunk_size - overlap
    return chunks


async def store_source_chunks(source_id: str, source_title: str, content: str, notebook_id: str):
    """Split source content into chunks and store in DB for retrieval"""
    await db.source_chunks.delete_many({"source_id": source_id})
    chunks = split_into_chunks(content)
    if not chunks:
        return
    docs = []
    for i, chunk_text in enumerate(chunks):
        docs.append({
            "source_id": source_id,
            "source_title": source_title,
            "notebook_id": notebook_id,
            "chunk_index": i,
            "text": chunk_text,
            "word_count": len(chunk_text.split())
        })
    if docs:
        await db.source_chunks.insert_many(docs)


def score_chunk_relevance(query: str, chunk_text: str) -> float:
    """Simple keyword-based relevance scoring"""
    query_words = set(query.lower().split())
    stop_words = {'the','a','an','is','are','was','were','be','been','being','have','has','had',
                  'do','does','did','will','would','could','should','may','might','shall','can',
                  'to','of','in','for','on','with','at','by','from','as','into','about','what',
                  'which','who','whom','this','that','these','those','am','or','and','but','not',
                  'no','all','any','my','your','his','her','its','our','their','i','you','he','she',
                  'it','we','they','me','him','us','them'}
    query_words = query_words - stop_words
    if not query_words:
        return 0.0
    chunk_lower = chunk_text.lower()
    matches = sum(1 for w in query_words if w in chunk_lower)
    return matches / len(query_words)


async def retrieve_relevant_chunks(query: str, notebook_id: str, source_ids: List[str] = None, top_k: int = 12) -> List[dict]:
    """Retrieve the most relevant chunks for a query using keyword scoring"""
    filter_q = {"notebook_id": notebook_id}
    if source_ids:
        filter_q["source_id"] = {"$in": source_ids}
    
    all_chunks = await db.source_chunks.find(filter_q, {"_id": 0}).to_list(500)
    
    if not all_chunks:
        return []
    
    for chunk in all_chunks:
        chunk["score"] = score_chunk_relevance(query, chunk["text"])
    
    all_chunks.sort(key=lambda c: c["score"], reverse=True)
    return all_chunks[:top_k]

async def generate_with_ai(prompt: str, system_message: str = "You are an expert research assistant.") -> str:
    """Generate content using GPT-5.1 via Emergent"""
    try:
        chat = LlmChat(
            api_key=EMERGENT_LLM_KEY,
            session_id=str(uuid.uuid4()),
            system_message=system_message
        ).with_model("openai", "gpt-5.1")
        
        user_message = UserMessage(text=prompt)
        response = await chat.send_message(user_message)
        return response
    except Exception as e:
        error_str = str(e)
        logger.error(f"AI generation error: {error_str}")
        if "Budget has been exceeded" in error_str or "budget" in error_str.lower():
            raise HTTPException(status_code=402, detail="LLM budget exceeded. Please go to Profile > Universal Key > Add Balance to top up.")
        raise HTTPException(status_code=500, detail=f"AI generation failed: {error_str}")

async def generate_slide_image(slide_title: str, slide_content: str, layout_type: str, theme: str, image_keyword: str = "") -> Optional[str]:
    """Generate a professional, contextually relevant image for a slide using AI"""
    try:
        # Use the image_keyword for more targeted image generation
        context = image_keyword if image_keyword else slide_title
        content_hint = slide_content[:200] if slide_content else context
        
        # Theme-specific colors
        theme_colors = {
            'tech': 'blue and white color scheme, technology aesthetic',
            'smart_home': 'green and teal color scheme, modern smart home aesthetic',
            'corporate': 'indigo and navy color scheme, corporate aesthetic',
            'finance': 'teal and dark green color scheme, financial aesthetic',
            'health': 'teal (#004D40) and warm gold (#C8A86B) color scheme, healthcare aesthetic',
            'education': 'amber and warm orange color scheme, educational aesthetic'
        }
        color_scheme = theme_colors.get(theme, 'professional blue color scheme')
        
        prompt = f"""Create a professional, photorealistic illustration for a business presentation slide.

Topic: {context}
Context: {content_hint}

Style requirements:
- {color_scheme}
- Clean, modern professional design suitable for a corporate presentation
- Photorealistic or high-quality 3D rendered style
- Subtle, relevant visual metaphor for the topic
- Light or white background that blends well with slide content
- No text, no labels, no words in the image
- Landscape orientation (wider than tall)
- Soft lighting, professional corporate feel
- High resolution, crisp detail"""

        logger.info(f"Generating image for slide: {slide_title} (keyword: {image_keyword})")
        
        image_gen = OpenAIImageGeneration(api_key=EMERGENT_LLM_KEY)
        images = await image_gen.generate_images(
            prompt=prompt,
            model="gpt-image-1",
            number_of_images=1
        )
        
        if images and len(images) > 0:
            # Compress: resize to 800x600 JPEG for smaller PPTX file size
            from PIL import Image as PILImage
            import io
            img = PILImage.open(io.BytesIO(images[0]))
            img = img.convert('RGB')
            img = img.resize((800, 600), PILImage.LANCZOS)
            buffer = io.BytesIO()
            img.save(buffer, format='JPEG', quality=80, optimize=True)
            compressed = buffer.getvalue()
            image_base64 = base64.b64encode(compressed).decode('utf-8')
            logger.info(f"Image for '{slide_title}': {len(images[0])}B -> {len(compressed)}B")
            return image_base64
        return None
    except Exception as e:
        logger.error(f"Image generation error for '{slide_title}': {e}")
        return None

# ─── ROUTES ─────────────────────────────────────────────────────────────────

@api_router.get("/")
async def root():
    return {"message": "NotebookLM API Ready"}

@api_router.post("/status", response_model=StatusCheck)
async def create_status_check(input: StatusCheckCreate):
    status_dict = input.model_dump()
    status_obj = StatusCheck(**status_dict)
    doc = status_obj.model_dump()
    doc['timestamp'] = doc['timestamp'].isoformat()
    _ = await db.status_checks.insert_one(doc)
    return status_obj

@api_router.get("/status", response_model=List[StatusCheck])
async def get_status_checks():
    status_checks = await db.status_checks.find({}, {"_id": 0}).to_list(1000)
    for check in status_checks:
        if isinstance(check['timestamp'], str):
            check['timestamp'] = datetime.fromisoformat(check['timestamp'])
    return status_checks

# ─── SOURCE MANAGEMENT ──────────────────────────────────────────────────────

@api_router.post("/sources/upload")
async def upload_source(file: UploadFile = File(...), notebook_id: str = "default"):
    """Upload and process a document file (PDF, DOCX, TXT, XLSX, PPTX, images)"""
    try:
        content = await file.read()
        filename = file.filename or "unknown"
        file_ext = filename.rsplit('.', 1)[-1].lower() if '.' in filename else 'txt'
        
        logger.info(f"Processing file: {filename}, extension: {file_ext}, size: {len(content)} bytes")
        
        # Extract text based on file type
        if file_ext == 'pdf':
            text = await extract_text_from_pdf(content)
            file_type = 'pdf'
        elif file_ext in ['docx', 'doc']:
            text = await extract_text_from_docx(content)
            file_type = 'doc'
        elif file_ext in ['pptx', 'ppt']:
            text = await extract_text_from_pptx(content)
            file_type = 'ppt'
        elif file_ext in ['xlsx', 'xls']:
            text = await extract_text_from_excel(content)
            file_type = 'xlsx'
        elif file_ext in ['png', 'jpg', 'jpeg', 'gif', 'webp', 'bmp', 'tiff', 'tif']:
            text = await extract_text_from_image(content, filename)
            file_type = 'image'
        elif file_ext == 'txt':
            text = content.decode('utf-8', errors='ignore')
            file_type = 'txt'
        else:
            # Fallback: try ZIP-based (DOCX/PPTX are ZIPs), then text
            try:
                if content[:4] == b'PK\x03\x04':
                    text = await extract_text_from_docx(content)
                    file_type = 'doc'
                else:
                    text = content.decode('utf-8', errors='ignore')
                    file_type = 'txt'
            except Exception:
                text = content.decode('utf-8', errors='ignore')
                file_type = 'txt'
        
        logger.info(f"Extracted text length: {len(text)} chars for type: {file_type}")
        
        if not text or len(text.strip()) < 5:
            raise HTTPException(status_code=400, detail=f"Could not extract readable content from file. File type: {file_ext}")
        
        # Create source document
        source = Source(
            title=filename.rsplit('.', 1)[0] if '.' in filename else filename,
            type=file_type,
            content=text,
            chunks=count_chunks(text),
            size=f"{len(content) / 1024:.1f} KB",
            status="indexed",
            notebook_id=notebook_id
        )
        
        # Store in MongoDB
        doc = source.model_dump()
        doc['created_at'] = doc['created_at'].isoformat()
        await db.sources.insert_one(doc)
        
        # Store chunks for retrieval
        await store_source_chunks(source.id, source.title, text, notebook_id)
        
        return {
            "id": source.id,
            "title": source.title,
            "type": source.type,
            "chunks": source.chunks,
            "size": source.size,
            "status": source.status,
            "content_preview": text[:200] + "..." if len(text) > 200 else text
        }
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Upload error: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@api_router.post("/sources/url")
async def add_url_source(data: SourceCreate):
    """Add a URL as a source"""
    try:
        if not data.url:
            raise HTTPException(status_code=400, detail="URL is required")
        
        text = await extract_text_from_url(data.url)
        if not text:
            raise HTTPException(status_code=400, detail="Could not extract text from URL")
        
        source = Source(
            title=data.title or data.url,
            type="url",
            content=text,
            chunks=count_chunks(text),
            size="—",
            status="indexed",
            url=data.url,
            notebook_id=data.notebook_id
        )
        
        doc = source.model_dump()
        doc['created_at'] = doc['created_at'].isoformat()
        await db.sources.insert_one(doc)
        
        # Store chunks for retrieval
        await store_source_chunks(source.id, source.title, text, data.notebook_id)
        
        return {
            "id": source.id,
            "title": source.title,
            "type": source.type,
            "chunks": source.chunks,
            "status": source.status,
            "url": source.url
        }
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"URL add error: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@api_router.get("/sources")
async def get_sources(notebook_id: str = "default"):
    """Get all sources for a notebook"""
    sources = await db.sources.find(
        {"notebook_id": notebook_id}, 
        {"_id": 0, "content": 0}  # Exclude large content field
    ).to_list(100)
    return sources

@api_router.get("/sources/{source_id}")
async def get_source(source_id: str):
    """Get a specific source with content"""
    source = await db.sources.find_one({"id": source_id}, {"_id": 0})
    if not source:
        raise HTTPException(status_code=404, detail="Source not found")
    return source

@api_router.delete("/sources/{source_id}")
async def delete_source(source_id: str):
    """Delete a source"""
    result = await db.sources.delete_one({"id": source_id})
    if result.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Source not found")
    return {"status": "deleted"}

# ─── AI GENERATION ──────────────────────────────────────────────────────────

@api_router.post("/generate", response_model=GenerateResponse)
async def generate_output(request: GenerateRequest):
    """Generate AI-powered output from sources"""
    
    # Get all indexed sources for the notebook
    sources = await db.sources.find(
        {"notebook_id": request.notebook_id, "status": "indexed"},
        {"_id": 0}
    ).to_list(50)
    
    if not sources:
        raise HTTPException(status_code=400, detail="No indexed sources found. Please add sources first.")
    
    # Combine source content (limit to avoid token overflow)
    combined_content = ""
    for src in sources:
        combined_content += f"\n\n--- SOURCE: {src['title']} ---\n{src.get('content', '')[:10000]}"
    combined_content = combined_content[:40000]  # Limit total content
    
    source_titles = [s['title'] for s in sources]
    title = request.title or f"{request.output_type.title()} - Research Summary"
    
    # Generate based on output type
    if request.output_type == "slides":
        return await generate_slides(combined_content, source_titles, title)
    elif request.output_type == "report":
        return await generate_report(combined_content, source_titles, title)
    elif request.output_type == "mindmap":
        return await generate_mindmap(combined_content, source_titles, title)
    elif request.output_type == "flashcards":
        return await generate_flashcards(combined_content, source_titles, title)
    elif request.output_type == "quiz":
        return await generate_quiz(combined_content, source_titles, title)
    elif request.output_type == "audio":
        return await generate_audio_script(combined_content, source_titles, title)
    elif request.output_type == "datatable":
        return await generate_datatable(combined_content, source_titles, title)
    elif request.output_type == "infographic":
        return await generate_infographic(combined_content, source_titles, title)
    else:
        raise HTTPException(status_code=400, detail=f"Unknown output type: {request.output_type}")

def detect_theme(content: str) -> str:
    """Detect appropriate theme from content"""
    content_lower = content.lower()
    health_kw = ['health', 'medical', 'clinical', 'patient', 'hospital', 'care', 'pharma', 'therapy']
    if any(k in content_lower for k in health_kw):
        return 'health'
    tech_kw = ['software', 'platform', 'api', 'cloud', 'data', 'algorithm', 'ai', 'machine learning']
    if any(k in content_lower for k in tech_kw):
        return 'tech'
    finance_kw = ['finance', 'revenue', 'budget', 'investment', 'cost', 'profit']
    if any(k in content_lower for k in finance_kw):
        return 'finance'
    return 'corporate'


async def generate_slides(content: str, sources: List[str], title: str) -> GenerateResponse:
    """Generate PowerPoint slide content using AI with rich formatting"""
    
    prompt = f"""Based on the following research content from {len(sources)} sources, create a professional PowerPoint presentation.

SOURCES ANALYZED: {', '.join(sources)}

CONTENT:
{content}

Create a presentation with exactly 8-10 slides. For each slide, provide:
1. A clear, concise title (max 8 words)
2. 3-5 bullet points with key insights (each bullet max 20 words)
3. Speaker notes (2-3 sentences explaining the slide)
4. Layout type — choose the BEST fit for the content:
   - "title" — ONLY for slide 1 (cover slide with subtitle)
   - "image-right" — Use for 2-3 key concept slides where a visual would enhance understanding. The image will be AI-generated based on the imageKeyword.
   - "image-left" — Use for 1-2 slides to add visual variety, alternating with image-right.
   - "two-column" — Use for comparisons or when content naturally splits into two groups
   - "timeline" — Use ONLY when content describes sequential phases/steps (max 4-5 steps)
   - "bullets" — Default for content-heavy slides that don't fit other layouts
   - "quote" — Use for a key takeaway or conclusion slide
5. imageKeyword: A specific 2-4 word phrase describing what the image should show (e.g., "hospital patient monitoring", "digital health dashboard", "team collaboration meeting"). This keyword drives the AI image generator — be specific and relevant.
6. An icon from: home, briefcase, chart, users, cog, shield, cloud, mobile, check, lightbulb, rocket, target, clock, globe, lock, server, database, code, team, growth
7. highlight: A key phrase to emphasize (optional, max 10 words)

LAYOUT DISTRIBUTION RULES:
- Slide 1 MUST be "title"
- Use "image-left" or "image-right" for 3-4 slides total (these will have AI-generated images)
- Use "two-column" or "timeline" for 1-2 slides
- Use "bullets" for 2-3 content-heavy slides
- End with "quote" or "bullets" for conclusion

Format your response as JSON array:
[
  {{
    "slideNumber": 1,
    "title": "Slide Title Here",
    "subtitle": "Descriptive subtitle for the presentation",
    "bullets": ["Bullet point 1", "Bullet point 2", "Bullet point 3"],
    "notes": "Speaker notes explaining this slide content.",
    "layout": "title",
    "imageKeyword": "topic relevant visual",
    "icon": "rocket",
    "highlight": ""
  }}
]

Only output the JSON array, no other text."""

    response = await generate_with_ai(prompt, "You are an expert presentation designer. Create visually engaging, professional slide content with varied layouts.")
    
    # Parse JSON response
    try:
        import json
        # Clean response - extract JSON from response
        response_clean = response.strip()
        if response_clean.startswith("```"):
            response_clean = response_clean.split("```")[1]
            if response_clean.startswith("json"):
                response_clean = response_clean[4:]
        slides_data = json.loads(response_clean)
        
        # Ensure all slides have required fields with defaults
        for slide in slides_data:
            slide.setdefault('layout', 'bullets')
            slide.setdefault('imageKeyword', 'business')
            slide.setdefault('icon', 'briefcase')
            slide.setdefault('subtitle', '')
            slide.setdefault('highlight', '')
            slide.setdefault('bullets', [])
            slide.setdefault('notes', '')
    except Exception as e:
        logger.error(f"JSON parse error: {e}")
        # Fallback if JSON parsing fails
        slides_data = [
            {"slideNumber": 1, "title": title, "subtitle": f"Based on {len(sources)} sources", "bullets": [], "notes": "Title slide", "layout": "title", "imageKeyword": "presentation", "icon": "briefcase"},
            {"slideNumber": 2, "title": "Key Findings", "bullets": ["Analysis in progress", "See full content below"], "notes": response[:500], "layout": "bullets", "imageKeyword": "analysis", "icon": "chart"}
        ]
    
    # Determine theme
    theme = detect_theme(content)
    
    # NOTE: Image generation is disabled to avoid gateway timeout (60s limit)
    # The presentation will be generated with placeholder icons instead
    # To enable images, run locally or increase gateway timeout
    logger.info(f"Slides generated: {len(slides_data)} (images disabled due to timeout constraints)")
    
    # Create text version for preview
    text_content = f"📊 PRESENTATION: {title}\n\nSources: {', '.join(sources)}\nTheme: {theme}\n\n"
    for slide in slides_data:
        text_content += f"\n{'='*50}\nSlide {slide.get('slideNumber', '?')}: {slide.get('title', 'Untitled')} [{slide.get('layout', 'bullets')}]\n{'='*50}\n"
        if slide.get('subtitle'):
            text_content += f"Subtitle: {slide.get('subtitle')}\n"
        for bullet in slide.get('bullets', []):
            text_content += f"• {bullet}\n"
        if slide.get('highlight'):
            text_content += f"💡 Highlight: {slide.get('highlight')}\n"
        text_content += f"🖼️ Image: {slide.get('imageKeyword', 'business')}"
        if slide.get('imageBase64'):
            text_content += " ✅ Generated"
        text_content += f"\n📝 Notes: {slide.get('notes', '')}\n"
    
    return GenerateResponse(
        id=str(uuid.uuid4()),
        type="slides",
        title=title,
        content=text_content,
        slides_data=slides_data,
        theme=theme
    )

async def generate_report(content: str, sources: List[str], title: str) -> GenerateResponse:
    """Generate a visual research report with structured sections"""
    
    prompt = f"""Based on the following research content from {len(sources)} sources, create a HIGHLY VISUAL research report. This report will be rendered as an infographic-style visual — prioritize images, workflows, and key stats over text.

SOURCES: {', '.join(sources)}

CONTENT:
{content}

Create a report with exactly 5 sections. Each section is IMAGE-FIRST with minimal text.

For each section provide:
1. title: Section heading (max 5 words)
2. subtitle: A single-sentence takeaway (max 10 words)
3. bullets: ONLY 2 short bullet points (max 8 words each) — key data only
4. stat: A key number/percentage from the content for this section (e.g., "85%", "24/7", "4 Phases", "3x Faster"). If no real stat, use "".
5. statLabel: Short label for the stat (max 4 words)
6. imageKeyword: Specific 3-5 word phrase for AI image (e.g., "remote patient monitoring dashboard", "clinical team hospital ward")
7. visualType: Choose the BEST visual for this section:
   - "workflow" — for processes with sequential steps (use for 2-3 sections)
   - "stat" — for sections with a strong numeric highlight (use for 1-2 sections)
   - "none" — only if no visual fits
8. workflowSteps: If visualType is "workflow", list 3-5 short step labels (max 4 words each)
9. icon: from chart, users, clock, target, shield, globe, lightbulb, rocket, cog, check, growth, home, briefcase

Format as JSON array:
[
  {{
    "sectionNumber": 1,
    "title": "Section Title",
    "subtitle": "One line takeaway here",
    "bullets": ["Point one here", "Point two here"],
    "stat": "85%",
    "statLabel": "Coverage Rate",
    "imageKeyword": "specific visual concept here",
    "visualType": "workflow",
    "workflowSteps": ["Step One", "Step Two", "Step Three"],
    "icon": "chart"
  }}
]

RULES:
- MINIMAL TEXT. 2 bullets max per section, 8 words max each.
- At least 2 sections MUST have visualType "workflow"
- At least 1 section MUST have visualType "stat" with a real number
- Every section MUST have a specific imageKeyword
- subtitle is required for every section
- Only output the JSON array, no other text."""

    response = await generate_with_ai(prompt, "You are a visual report designer. Create image-heavy, minimal-text report structures.")
    
    # Parse JSON
    try:
        import json
        clean = response.strip()
        if clean.startswith('```'): clean = clean.split('\n', 1)[1].rsplit('```', 1)[0]
        report_sections = json.loads(clean)
        for sec in report_sections:
            sec.setdefault('title', 'Section')
            sec.setdefault('subtitle', '')
            sec.setdefault('bullets', [])
            sec['bullets'] = sec['bullets'][:2]  # Hard cap at 2 bullets
            sec.setdefault('stat', '')
            sec.setdefault('statLabel', '')
            sec.setdefault('imageKeyword', sec.get('title', ''))
            sec.setdefault('visualType', 'none')
            sec.setdefault('workflowSteps', [])
            sec.setdefault('icon', 'briefcase')
    except Exception as e:
        logger.error(f"Report JSON parse error: {e}")
        report_sections = [
            {"sectionNumber": 1, "title": title, "bullets": ["Report generated from source documents"], "imageKeyword": "research analysis", "visualType": "none", "workflowSteps": [], "icon": "chart"}
        ]
    
    # Generate minimal text summary for preview
    text_content = f"VISUAL REPORT: {title}\nGenerated: {datetime.now().strftime('%B %d, %Y')}\nSources: {len(sources)} documents\n\n"
    for sec in report_sections:
        text_content += f"## {sec['title']}\n"
        if sec.get('subtitle'):
            text_content += f"{sec['subtitle']}\n"
        if sec.get('stat'):
            text_content += f"[{sec['stat']}] {sec.get('statLabel', '')}\n"
        for b in sec.get('bullets', [])[:2]:
            text_content += f"- {b}\n"
        text_content += "\n"
    
    return GenerateResponse(
        id=str(uuid.uuid4()),
        type="report",
        title=title,
        content=text_content,
        slides_data=report_sections,
        theme=detect_theme(content)
    )

async def generate_mindmap(content: str, sources: List[str], title: str) -> GenerateResponse:
    """Generate a structured mind map"""
    import json
    prompt = f"""Based on the following content from {len(sources)} sources, create a mind map structure.

SOURCES: {', '.join(sources)}
CONTENT:
{content[:5000]}

Return a JSON object with this exact structure:
{{
  "center": "Central Topic (max 5 words)",
  "branches": [
    {{
      "label": "Branch Name (max 4 words)",
      "color": "teal",
      "children": [
        {{"label": "Sub-topic 1 (max 5 words)"}},
        {{"label": "Sub-topic 2 (max 5 words)"}},
        {{"label": "Sub-topic 3 (max 5 words)"}}
      ]
    }}
  ]
}}

RULES:
- Exactly 5-6 main branches
- Each branch has 3-4 children
- color must be one of: teal, gold, blue, red, purple, green
- Labels are SHORT (max 5 words)
- Only output JSON, no other text."""

    response = await generate_with_ai(prompt, "Mind map designer. Concise labels.")
    try:
        clean = response.strip()
        if clean.startswith('```'): clean = clean.split('\n', 1)[1].rsplit('```', 1)[0]
        mm_data = json.loads(clean)
    except:
        mm_data = {"center": title, "branches": [
            {"label": "Core Concepts", "color": "teal", "children": [{"label": "Overview"}, {"label": "Key Terms"}]},
            {"label": "Key Findings", "color": "blue", "children": [{"label": "Data Points"}, {"label": "Analysis"}]},
            {"label": "Applications", "color": "green", "children": [{"label": "Use Cases"}, {"label": "Impact"}]},
        ]}
    
    text_preview = f"MIND MAP: {title}\n\nCenter: {mm_data.get('center','')}\n"
    for b in mm_data.get('branches', []):
        text_preview += f"\n{b['label']}:\n"
        for c in b.get('children', []):
            text_preview += f"  - {c['label']}\n"
    
    return GenerateResponse(id=str(uuid.uuid4()), type="mindmap", title=title, content=text_preview, slides_data=mm_data)

async def generate_flashcards(content: str, sources: List[str], title: str) -> GenerateResponse:
    """Generate structured flashcards"""
    import json
    prompt = f"""Based on the following content from {len(sources)} sources, create 12 study flashcards.

SOURCES: {', '.join(sources)}
CONTENT:
{content[:5000]}

Return a JSON array of flashcard objects:
[
  {{
    "id": 1,
    "question": "Clear question (max 15 words)",
    "answer": "Concise answer (max 25 words)",
    "category": "Category Name",
    "difficulty": "Easy"
  }}
]

RULES:
- Exactly 12 cards
- difficulty: "Easy", "Medium", or "Hard" (mix of all three)
- category: group related cards (3-4 unique categories)
- Questions test understanding, not memorization
- Only output JSON array."""

    response = await generate_with_ai(prompt, "Expert educator. Create flashcards.")
    try:
        clean = response.strip()
        if clean.startswith('```'): clean = clean.split('\n', 1)[1].rsplit('```', 1)[0]
        cards = json.loads(clean)
        for i, c in enumerate(cards):
            c.setdefault('id', i + 1)
            c.setdefault('difficulty', 'Medium')
            c.setdefault('category', 'General')
    except:
        cards = [{"id": i+1, "question": f"Question {i+1}", "answer": "Answer pending", "category": "General", "difficulty": "Medium"} for i in range(6)]
    
    text_preview = f"FLASHCARDS: {title}\n{len(cards)} cards\n\n"
    for c in cards:
        text_preview += f"Q{c['id']}: {c['question']}\nA: {c['answer']}\n\n"
    
    return GenerateResponse(id=str(uuid.uuid4()), type="flashcards", title=title, content=text_preview, slides_data=cards)

async def generate_quiz(content: str, sources: List[str], title: str) -> GenerateResponse:
    """Generate a structured quiz"""
    import json
    prompt = f"""Based on the following content from {len(sources)} sources, create a quiz.

SOURCES: {', '.join(sources)}
CONTENT:
{content[:5000]}

Return a JSON object:
{{
  "mcq": [
    {{"id": 1, "question": "Question text?", "options": ["A) Option", "B) Option", "C) Option", "D) Option"], "correct": "B", "explanation": "Brief explanation"}}
  ],
  "truefalse": [
    {{"id": 1, "statement": "Statement text.", "answer": true, "explanation": "Brief explanation"}}
  ],
  "short": [
    {{"id": 1, "question": "Question text?", "sampleAnswer": "Sample answer text"}}
  ]
}}

RULES:
- 5 MCQ questions, 4 TF questions, 3 short answer questions
- Questions must be based on the source content
- Only output JSON."""

    response = await generate_with_ai(prompt, "Expert test creator.")
    try:
        clean = response.strip()
        if clean.startswith('```'): clean = clean.split('\n', 1)[1].rsplit('```', 1)[0]
        quiz_data = json.loads(clean)
    except:
        quiz_data = {"mcq": [], "truefalse": [], "short": []}
    
    quiz_data.setdefault('mcq', [])
    quiz_data.setdefault('truefalse', [])
    quiz_data.setdefault('short', [])
    total = len(quiz_data['mcq']) + len(quiz_data['truefalse']) + len(quiz_data['short'])
    
    text_preview = f"QUIZ: {title}\n{total} questions\n\n"
    for q in quiz_data.get('mcq', []):
        text_preview += f"MCQ: {q['question']}\n"
    for q in quiz_data.get('truefalse', []):
        text_preview += f"T/F: {q['statement']}\n"
    for q in quiz_data.get('short', []):
        text_preview += f"Short: {q['question']}\n"
    
    return GenerateResponse(id=str(uuid.uuid4()), type="quiz", title=title, content=text_preview, slides_data=quiz_data)

async def generate_audio_script(content: str, sources: List[str], title: str) -> GenerateResponse:
    """Generate a podcast-style audio script"""
    
    prompt = f"""Based on the following content from {len(sources)} sources, create an engaging podcast script for two hosts.

SOURCES: {', '.join(sources)}

CONTENT:
{content}

Create a ~10 minute podcast script with:
- Introduction and hook
- Discussion of key topics
- Interesting insights and analysis
- Conclusion and takeaways

Format as a dialogue between HOST A and HOST B.
Make it conversational, engaging, and informative."""

    response = await generate_with_ai(prompt, "You are an expert podcast scriptwriter. Create engaging audio content.")
    
    script = f"""🎙️ PODCAST SCRIPT: {title}

Duration: ~10-12 minutes
Sources: {', '.join(sources)}

{'═'*50}

{response}

{'═'*50}

[End of script - Ready for recording]
"""
    
    return GenerateResponse(
        id=str(uuid.uuid4()),
        type="audio",
        title=title,
        content=script
    )

async def generate_datatable(content: str, sources: List[str], title: str) -> GenerateResponse:
    """Generate structured data tables"""
    import json
    prompt = f"""Based on the following content from {len(sources)} sources, extract data into structured tables.

SOURCES: {', '.join(sources)}
CONTENT:
{content[:5000]}

Return a JSON object:
{{
  "tables": [
    {{
      "title": "Table Name",
      "headers": ["Column 1", "Column 2", "Column 3"],
      "rows": [
        ["Value 1", "Value 2", "Value 3"],
        ["Value 4", "Value 5", "Value 6"]
      ]
    }}
  ],
  "stats": [
    {{"label": "Metric Name", "value": "42", "description": "Brief context"}}
  ]
}}

RULES:
- Create 2-3 tables with real data from the sources
- Each table has 3-5 columns and 4-8 rows
- Extract 3-5 key statistics/metrics
- Values must come from source content
- Only output JSON."""

    response = await generate_with_ai(prompt, "Expert data analyst.")
    try:
        clean = response.strip()
        if clean.startswith('```'): clean = clean.split('\n', 1)[1].rsplit('```', 1)[0]
        dt_data = json.loads(clean)
    except:
        dt_data = {"tables": [{"title": "Data Summary", "headers": ["Source", "Type", "Key Finding"], "rows": [[s, "Document", "Pending extraction"] for s in sources]}], "stats": []}
    
    dt_data.setdefault('tables', [])
    dt_data.setdefault('stats', [])
    
    text_preview = f"DATA TABLE: {title}\n{len(dt_data['tables'])} tables, {len(dt_data['stats'])} stats\n\n"
    for t in dt_data['tables']:
        text_preview += f"{t['title']}: {len(t.get('rows',[]))} rows\n"
    
    return GenerateResponse(id=str(uuid.uuid4()), type="datatable", title=title, content=text_preview, slides_data=dt_data)

async def generate_infographic(content: str, sources: List[str], title: str) -> GenerateResponse:
    """Generate a HIGHLY VISUAL infographic — image-first, minimal text"""
    
    prompt = f"""Based on the following content from {len(sources)} sources, create a HIGHLY VISUAL infographic. This will be rendered as an image-first visual — prioritize visuals, workflows, and stats over text.

SOURCES: {', '.join(sources)}

CONTENT:
{content}

Create exactly 6 sections. Each section is IMAGE-DOMINANT with minimal text.

For each section provide:
1. title: Section heading (max 5 words)
2. subtitle: Single-sentence takeaway (max 10 words)
3. bullets: ONLY 2 short bullet points (max 8 words each) — key data only, no sentences
4. stat: A key number/metric from the content (e.g., "85%", "24/7", "5 Phases", "3x"). Use "" if none.
5. statLabel: Short label for stat (max 4 words). Use "" if no stat.
6. imageKeyword: Specific 3-5 word phrase for AI image (e.g., "remote patient monitoring dashboard", "clinical team hospital ward")
7. visualType: Choose the BEST visual:
   - "workflow" — for processes/flows with steps (use for 2-3 sections)
   - "stat" — for sections with numeric highlights (use for 1-2 sections)
   - "none" — only if nothing fits
8. workflowSteps: If visualType is "workflow", list 3-5 step labels (max 4 words each)
9. icon: from chart, users, clock, target, shield, globe, lightbulb, rocket, cog, check, growth, home, briefcase

Format as JSON array:
[
  {{
    "sectionNumber": 1,
    "title": "Section Title",
    "subtitle": "One line takeaway",
    "bullets": ["Short point one", "Short point two"],
    "stat": "5",
    "statLabel": "Key Modules",
    "imageKeyword": "specific visual concept phrase",
    "visualType": "workflow",
    "workflowSteps": ["Step One", "Step Two", "Step Three"],
    "icon": "chart"
  }}
]

RULES:
- MINIMAL TEXT. Max 2 bullets per section, max 8 words each. No long sentences.
- At least 2 sections MUST have visualType "workflow" with workflowSteps
- At least 1 section MUST have visualType "stat" with a real number
- Every section MUST have subtitle and imageKeyword
- Only output the JSON array, no other text."""

    response = await generate_with_ai(prompt, "You are a visual infographic designer. Create image-heavy, minimal-text infographic structures.")
    
    # Parse JSON response
    try:
        import json
        response_clean = response.strip()
        if response_clean.startswith("```"):
            response_clean = response_clean.split("```")[1]
            if response_clean.startswith("json"):
                response_clean = response_clean[4:]
        infographic_data = json.loads(response_clean)
        
        # Ensure all sections have required fields + enforce limits
        for section in infographic_data:
            section.setdefault('visualType', 'none')
            section.setdefault('icon', 'briefcase')
            section.setdefault('stat', '')
            section.setdefault('statLabel', '')
            section.setdefault('subtitle', '')
            section.setdefault('bullets', [])
            section['bullets'] = section['bullets'][:2]  # Hard cap at 2 bullets
            section.setdefault('workflowSteps', [])
            section.setdefault('color', 'blue')
            section.setdefault('imageKeyword', section.get('title', ''))
    except Exception as e:
        logger.error(f"Infographic JSON parse error: {e}")
        infographic_data = [
            {"sectionNumber": 1, "title": title, "stat": str(len(sources)), "statLabel": "Sources Analyzed", "bullets": [], "icon": "chart", "visualType": "none", "color": "blue"},
            {"sectionNumber": 2, "title": "Key Findings", "bullets": ["Analysis generated from source documents"], "icon": "lightbulb", "visualType": "none", "color": "green"}
        ]
    
    # Determine theme based on content
    content_lower = content.lower()
    if any(word in content_lower for word in ['tech', 'software', 'code', 'development', 'app', 'digital']):
        theme = 'tech'
    elif any(word in content_lower for word in ['health', 'medical', 'patient', 'care', 'hospital']):
        theme = 'health'
    elif any(word in content_lower for word in ['finance', 'money', 'budget', 'cost', 'revenue']):
        theme = 'finance'
    else:
        theme = 'corporate'
    
    # Generate header image for the infographic
    logger.info("Generating infographic header image...")
    header_image = None
    try:
        header_image = await generate_slide_image(
            title,
            ' '.join([s.get('title', '') for s in infographic_data[:3]]),
            'infographic',
            theme
        )
        if header_image:
            logger.info("Header image generated successfully")
    except Exception as e:
        logger.error(f"Header image generation failed: {e}")
    
    # Create text preview
    text_content = f"""📊 INFOGRAPHIC: {title}

Sources: {', '.join(sources)}
Theme: {theme}
Sections: {len(infographic_data)}

{'═'*60}
"""
    for section in infographic_data:
        text_content += f"\n[{section.get('visualType', 'list').upper()}] {section.get('title', 'Section')}\n"
        if section.get('stat'):
            text_content += f"   📈 {section.get('stat')} - {section.get('statLabel', '')}\n"
        for bullet in section.get('bullets', []):
            text_content += f"   • {bullet}\n"
        text_content += f"   🎨 Icon: {section.get('icon')} | Color: {section.get('color')}\n"
    
    text_content += f"\n{'═'*60}\n"
    
    # Add header image to first section if available
    if header_image and infographic_data:
        infographic_data[0]['headerImage'] = header_image
    
    return GenerateResponse(
        id=str(uuid.uuid4()),
        type="infographic",
        title=title,
        content=text_content,
        slides_data=infographic_data,  # Reuse slides_data field for infographic sections
        theme=theme
    )

# ─── OUTPUT STORAGE ─────────────────────────────────────────────────────────

class SaveOutputRequest(BaseModel):
    id: str
    type: str
    title: str
    content: str
    slides_data: Optional[Union[List[dict], dict]] = None  # List for slides/report/infographic/flashcards, dict for mindmap/quiz/datatable
    notebook_id: str = "default"

@api_router.post("/outputs")
async def save_output(request: SaveOutputRequest):
    """Save a generated output to the database"""
    doc = {
        "id": request.id,
        "type": request.type,
        "title": request.title,
        "content": request.content,
        "slides_data": request.slides_data,
        "notebook_id": request.notebook_id,
        "created_at": datetime.now(timezone.utc).isoformat(),
        "size": f"{len(request.content) / 1024:.1f} KB"
    }
    await db.outputs.insert_one(doc)
    return {"status": "saved", "id": request.id}

@api_router.get("/outputs")
async def get_outputs(notebook_id: str = "default"):
    """Get all outputs for a notebook"""
    outputs = await db.outputs.find(
        {"notebook_id": notebook_id},
        {"_id": 0}
    ).sort("created_at", -1).to_list(100)
    return outputs

@api_router.delete("/outputs/{output_id}")
async def delete_output(output_id: str):
    """Delete an output"""
    result = await db.outputs.delete_one({"id": output_id})
    if result.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Output not found")
    return {"status": "deleted"}

# ─── SINGLE IMAGE GENERATION ────────────────────────────────────────────────

class GenerateImageRequest(BaseModel):
    slide_title: str
    slide_content: str
    layout: str = "bullets"
    theme: str = "corporate"
    image_keyword: str = ""

@api_router.post("/generate-image")
async def generate_single_image(request: GenerateImageRequest):
    """Generate a single image for a slide (for async image generation)"""
    try:
        image_base64 = await generate_slide_image(
            request.slide_title,
            request.slide_content,
            request.layout,
            request.theme,
            request.image_keyword
        )
        if image_base64:
            return {"success": True, "imageBase64": image_base64}
        return {"success": False, "error": "Image generation failed"}
    except Exception as e:
        logger.error(f"Single image generation error: {e}")
        return {"success": False, "error": str(e)}

# ─── CHAT WITH SOURCES ──────────────────────────────────────────────────────

class ChatMessage(BaseModel):
    role: str
    content: str

class ChatRequest(BaseModel):
    message: str
    notebook_id: str = "default"
    history: List[ChatMessage] = []
    depth: str = "balanced"
    source_ids: List[str] = []  # Empty = all sources
    mode: str = "general"  # general, compare, deep_analysis, executive_summary, qa_prep

class ChatResponse(BaseModel):
    response: str
    citations: List[dict] = []  # [{source: title, excerpt: text}]
    follow_ups: List[str] = []  # Suggested next questions

@api_router.post("/chat", response_model=ChatResponse)
async def chat_with_sources(request: ChatRequest):
    """Enhanced chat with RAG retrieval, modes, source filtering, and follow-up suggestions"""
    
    # Get indexed sources (optionally filtered)
    filter_q = {"notebook_id": request.notebook_id, "status": "indexed"}
    if request.source_ids:
        filter_q["id"] = {"$in": request.source_ids}
    
    sources = await db.sources.find(filter_q, {"_id": 0}).to_list(50)
    
    if not sources:
        return ChatResponse(
            response="No sources have been indexed yet. Please add some documents or URLs first.",
            citations=[], follow_ups=[]
        )
    
    source_titles = [s['title'] for s in sources]
    source_ids_list = [s['id'] for s in sources]
    
    # ── Phase A: Smart Retrieval ──
    # Try chunk-based retrieval first (RAG)
    relevant_chunks = await retrieve_relevant_chunks(
        request.message, request.notebook_id,
        source_ids=source_ids_list if request.source_ids else None,
        top_k=15
    )
    
    if relevant_chunks and any(c['score'] > 0.1 for c in relevant_chunks):
        # Use relevant chunks — group by source for context
        context_parts = []
        seen_sources = set()
        for chunk in relevant_chunks:
            if chunk['score'] < 0.05:
                continue
            src_label = chunk['source_title']
            if src_label not in seen_sources:
                context_parts.append(f"\n--- SOURCE: {src_label} ---")
                seen_sources.add(src_label)
            context_parts.append(chunk['text'])
        combined_content = '\n'.join(context_parts)
        retrieval_method = "semantic"
    else:
        # Fallback: send full source content with higher limits
        combined_content = ""
        for src in sources:
            combined_content += f"\n\n--- SOURCE: {src['title']} ---\n{src.get('content', '')}"
        combined_content = combined_content[:40000]
        retrieval_method = "full"
    
    # Detect translation requests
    msg_lower = request.message.lower()
    is_translation = any(kw in msg_lower for kw in ['translate', 'ترجم', 'translation', 'ترجمة'])
    
    # Build conversation history
    history_text = ""
    if request.history:
        recent = request.history[-12:]
        history_text = "\n\nCONVERSATION HISTORY:\n"
        for msg in recent:
            role_label = "USER" if msg.role == "user" else "ASSISTANT"
            history_text += f"{role_label}: {msg.content[:600]}\n"
    
    # ── Phase B: Chat Modes ──
    mode_instructions = {
        "general": "Answer the user's question thoroughly based on the sources. Be helpful and accurate.",
        "compare": "Compare and contrast information across the different sources. Highlight similarities, differences, contradictions, and complementary points. Structure your response with clear comparison categories.",
        "deep_analysis": "Provide an exhaustive, in-depth analysis. Cover every angle: root causes, implications, trends, risks, opportunities. Cross-reference between sources. Think critically and provide original insights.",
        "executive_summary": "Provide a concise executive summary suitable for senior leadership. Focus on key decisions, metrics, risks, and recommendations. Use bullet points and keep it actionable. Maximum 300 words.",
        "qa_prep": "Generate a thorough Q&A preparation document. Anticipate likely questions about this topic, provide detailed answers for each, and note potential follow-ups. Think like a subject matter expert preparing for a review meeting."
    }
    mode_note = mode_instructions.get(request.mode, mode_instructions["general"])
    
    # Depth instructions
    depth_instructions = {
        "fast": "Be concise and direct. Give a short, focused answer in 2-4 sentences. Skip unnecessary detail.",
        "balanced": "Provide a clear, well-structured answer with moderate detail. Use bullet points or short paragraphs.",
        "deep": "Provide a thorough, comprehensive answer with detailed analysis, examples, and cross-references between sources. Be as detailed as possible."
    }
    depth_note = depth_instructions.get(request.depth, depth_instructions["balanced"])
    
    source_list_str = ', '.join(source_titles[:10])
    source_filter_note = f"(User selected {len(source_titles)} specific sources)" if request.source_ids else f"(All {len(source_titles)} sources)"
    
    prompt = f"""You are a powerful research assistant with full access to the user's source documents.

AVAILABLE SOURCES {source_filter_note}: {source_list_str}

SOURCE CONTENT (retrieved via {retrieval_method} search):
{combined_content}
{history_text}

USER REQUEST: {request.message}

CHAT MODE: {mode_note}
RESPONSE DEPTH: {depth_note}

CAPABILITIES — You can and should handle ALL of the following when asked:
1. **Translation**: Translate source content to ANY language the user requests (Arabic, English, French, Spanish, etc.). Translate accurately and naturally, preserving meaning and formatting.
2. **Summarization**: Summarize individual sources or all sources together.
3. **Analysis & Comparison**: Compare sources, identify patterns, contradictions, or gaps.
4. **Extraction**: Pull out specific data, statistics, key terms, or quotes from sources.
5. **Q&A**: Answer factual questions grounded in the source content.
6. **Reformatting**: Convert content into tables, bullet points, numbered lists, or other formats.
7. **Writing**: Draft new content (emails, reports, briefs) based on source material.

CITATION RULES:
- When referencing information, include [Source: source_name] inline.
- Cite the specific source name when quoting or paraphrasing.
- If information is not in the sources, say so clearly.

FORMATTING RULES:
- Use markdown: **bold**, *italic*, bullet points, numbered lists, ### headers.
- Maintain conversation continuity — reference prior messages when relevant.

FOLLOW-UP: After your response, add a line "---FOLLOWUPS---" and then provide exactly 3 suggested follow-up questions the user might want to ask next, one per line. Make them specific and useful based on the conversation context. Do NOT number them."""

    system_msg = "You are an expert multilingual research assistant. You translate, analyze, summarize, compare, and answer questions about source documents. You support all languages."
    
    try:
        response = await generate_with_ai(prompt, system_msg)
    except HTTPException as he:
        error_msg = he.detail if he.detail else "AI generation failed"
        if he.status_code == 402:
            return ChatResponse(
                response="**LLM Budget Exceeded**\n\nYour Universal Key balance has run out. Please go to **Profile > Universal Key > Add Balance** to top up, or enable **Auto Top-Up** so this doesn't happen again.",
                citations=[], follow_ups=[]
            )
        return ChatResponse(
            response=f"An error occurred: {error_msg}\n\nPlease try again in a moment.",
            citations=[], follow_ups=[]
        )
    
    # Parse follow-ups from response
    follow_ups = []
    main_response = response
    if "---FOLLOWUPS---" in response:
        parts = response.split("---FOLLOWUPS---", 1)
        main_response = parts[0].strip()
        fu_text = parts[1].strip()
        follow_ups = [line.strip().lstrip('- ').lstrip('•').strip() for line in fu_text.split('\n') if line.strip() and len(line.strip()) > 5][:3]
    
    # Extract rich citations
    import re
    citations = []
    seen_cites = set()
    cite_pattern = re.findall(r'\[Source:\s*([^\]]+)\]', main_response)
    for cite_name in cite_pattern:
        clean = cite_name.strip()
        if clean not in seen_cites:
            seen_cites.add(clean)
            # Find a relevant excerpt from the chunk data
            excerpt = ""
            for chunk in relevant_chunks[:5]:
                if clean.lower() in chunk.get('source_title', '').lower():
                    excerpt = chunk['text'][:150] + "..."
                    break
            citations.append({"source": clean, "excerpt": excerpt})
    
    # Fallback: if no inline citations, add source names
    if not citations:
        for src in sources[:3]:
            if src['title'].lower() in main_response.lower():
                citations.append({"source": src['title'], "excerpt": ""})
    
    return ChatResponse(
        response=main_response,
        citations=citations,
        follow_ups=follow_ups
    )


@api_router.post("/chat/backfill-chunks")
async def backfill_chunks(notebook_id: str = "default"):
    """Backfill chunks for existing sources that don't have them yet"""
    sources = await db.sources.find(
        {"notebook_id": notebook_id, "status": "indexed"},
        {"_id": 0, "id": 1, "title": 1, "content": 1}
    ).to_list(100)
    count = 0
    for src in sources:
        existing = await db.source_chunks.count_documents({"source_id": src['id']})
        if existing == 0 and src.get('content'):
            await store_source_chunks(src['id'], src['title'], src['content'], notebook_id)
            count += 1
    return {"backfilled": count, "total_sources": len(sources)}


class ChatExportRequest(BaseModel):
    messages: List[dict]
    title: str = "Chat Export"

@api_router.post("/export/chat-docx")
async def export_chat_docx(req: ChatExportRequest):
    """Export chat as a professionally formatted Word document"""
    from docx import Document
    from docx.shared import Pt, Cm, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT
    from docx.oxml.ns import nsdecls
    from docx.oxml import parse_xml

    doc = Document()
    for section in doc.sections:
        section.top_margin = Cm(2)
        section.bottom_margin = Cm(2)
        section.left_margin = Cm(2.5)
        section.right_margin = Cm(2.5)

    style = doc.styles['Normal']
    style.font.name = 'Calibri'
    style.font.size = Pt(11)
    style.font.color.rgb = RGBColor(55, 65, 81)
    style.paragraph_format.line_spacing = 1.4

    h1 = doc.styles['Heading 1']
    h1.font.name = 'Calibri'
    h1.font.size = Pt(18)
    h1.font.bold = True
    h1.font.color.rgb = RGBColor(0, 77, 64)
    h1.paragraph_format.space_before = Pt(18)
    h1.paragraph_format.space_after = Pt(8)

    # Title bar
    tbl = doc.add_table(rows=1, cols=1)
    tbl.alignment = WD_TABLE_ALIGNMENT.CENTER
    cell = tbl.cell(0, 0)
    shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="004D40"/>')
    cell.paragraphs[0]._element.get_or_add_pPr().append(shading)
    cell_p = cell.paragraphs[0]
    cell_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    cell_p.space_before = Pt(16)
    cell_p.space_after = Pt(16)
    run = cell_p.add_run(req.title)
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = RGBColor(200, 168, 107)
    run.font.name = 'Calibri'

    doc.add_paragraph().space_after = Pt(4)
    meta = doc.add_paragraph()
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = meta.add_run(f'{len(req.messages)} messages  |  {datetime.now().strftime("%B %d, %Y")}  |  Knowledge Base')
    run.font.size = Pt(10)
    run.font.color.rgb = RGBColor(148, 163, 184)

    doc.add_paragraph().space_after = Pt(12)

    # Messages
    for msg in req.messages:
        role = msg.get('role', 'user')
        content = msg.get('content', '')
        time_str = msg.get('time', '')

        # Role header
        p = doc.add_paragraph()
        p.space_before = Pt(12)
        p.space_after = Pt(2)
        if role == 'user':
            run = p.add_run('You')
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 77, 64)
        else:
            run = p.add_run('AI Assistant')
            run.font.bold = True
            run.font.color.rgb = RGBColor(14, 165, 233)
        if time_str:
            run = p.add_run(f'  {time_str}')
            run.font.size = Pt(9)
            run.font.color.rgb = RGBColor(148, 163, 184)

        # Content
        cp = doc.add_paragraph()
        cp.space_after = Pt(4)
        # Simple markdown stripping for clean text
        import re as _re
        clean_content = _re.sub(r'\*\*([^*]+)\*\*', r'\1', content)
        clean_content = _re.sub(r'\*([^*]+)\*', r'\1', clean_content)
        run = cp.add_run(clean_content)
        run.font.size = Pt(11)

        # Citations
        cites = msg.get('citations', [])
        if cites:
            cp2 = doc.add_paragraph()
            cp2.space_after = Pt(2)
            run = cp2.add_run('Sources: ')
            run.font.size = Pt(9)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 77, 64)
            cite_names = [c.get('source', c) if isinstance(c, dict) else str(c) for c in cites]
            run = cp2.add_run(', '.join(cite_names))
            run.font.size = Pt(9)
            run.font.color.rgb = RGBColor(100, 116, 139)

        # Separator
        sep = doc.add_paragraph()
        sep.space_after = Pt(2)
        run = sep.add_run('\u2500' * 60)
        run.font.size = Pt(6)
        run.font.color.rgb = RGBColor(229, 231, 235)

    # Footer
    doc.add_paragraph().space_after = Pt(20)
    fp = doc.add_paragraph()
    fp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = fp.add_run('Generated by Knowledge Base  |  AI-Powered Research Assistant')
    run.font.size = Pt(9)
    run.font.color.rgb = RGBColor(148, 163, 184)

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    import re as _re
    safe_title = _re.sub(r'[^\x20-\x7E]', '', req.title).replace(' ', '_') or 'chat'
    filename = f'{safe_title}_chat.docx'
    return StreamingResponse(buf, media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                             headers={"Content-Disposition": f'attachment; filename="{filename}"'})


# ─── RFP GENERATOR ──────────────────────────────────────────────────────────

class RFPGenerateRequest(BaseModel):
    template_sections: List[str]
    project_name: str = ""
    client_name: str = ""
    tone: str = "Formal"
    additional_context: str = ""
    notebook_id: str = "default"

@api_router.post("/rfp/parse-template")
async def parse_rfp_template(file: UploadFile = File(...)):
    """Parse an uploaded RFP template and extract section headings"""
    content = await file.read()
    filename = file.filename.lower()
    
    text = ""
    if filename.endswith('.pdf'):
        text = await extract_text_from_pdf(content)
    elif filename.endswith('.docx'):
        text = await extract_text_from_docx(content)
    elif filename.endswith('.txt'):
        text = content.decode('utf-8', errors='ignore')
    else:
        raise HTTPException(status_code=400, detail="Unsupported file type. Use PDF, DOCX, or TXT.")
    
    if not text.strip():
        raise HTTPException(status_code=400, detail="Could not extract text from the template.")
    
    # Use AI to extract section headings from the template
    prompt = f"""Analyze this RFP template and extract ALL section headings/titles. Return them as a JSON array of strings.

TEMPLATE TEXT:
{text[:8000]}

Return ONLY a JSON array of section title strings, e.g.:
["Executive Summary", "Background & Context", "Scope of Work", ...]

If you cannot identify clear sections, return the standard enterprise RFP sections:
["Executive Summary", "Background & Context", "Definitions & Acronyms", "Project Objectives", "Scope of Work", "Technical Requirements", "Deliverables & Acceptance Criteria", "Timeline & Milestones", "Commercial Model & Budget", "Evaluation Criteria", "Instructions to Bidders", "Submission Requirements", "Legal Terms & Conditions", "Data Protection & Security Compliance", "Appendices"]

Only output the JSON array."""

    response = await generate_with_ai(prompt, "You are an expert at analyzing RFP document structures.")
    
    import json
    try:
        clean = response.strip()
        if clean.startswith('```'): clean = clean.split('\n', 1)[1].rsplit('```', 1)[0]
        sections = json.loads(clean)
        if not isinstance(sections, list):
            sections = []
    except:
        sections = ["Executive Summary", "Background & Context", "Definitions & Acronyms", "Project Objectives", "Scope of Work", "Technical Requirements", "Deliverables & Acceptance Criteria", "Timeline & Milestones", "Commercial Model & Budget", "Evaluation Criteria", "Instructions to Bidders", "Submission Requirements", "Legal Terms & Conditions", "Data Protection & Security Compliance", "Appendices"]
    
    return {"sections": sections, "template_text": text[:3000]}


@api_router.post("/rfp/generate")
async def generate_rfp(req: RFPGenerateRequest):
    """Generate enterprise-grade RFP content for a batch of sections"""
    import json
    sources_cursor = db.sources.find(
        {"notebook_id": req.notebook_id, "status": "indexed"},
        {"_id": 0, "title": 1, "content": 1}
    )
    sources = await sources_cursor.to_list(50)
    
    if not sources:
        raise HTTPException(status_code=400, detail="No indexed sources found.")
    
    combined = ""
    source_names = []
    for s in sources:
        source_names.append(s['title'])
        combined += f"\n\n[{s['title']}]:\n{s['content'][:2500]}"
    combined = combined[:8000]
    
    tone_map = {"Formal": "formal bureaucratic procurement", "Technical": "technical specifications-focused", "Executive": "executive strategic decision-maker", "Proposal-style": "persuasive benefits-focused"}
    sections_list = "\n".join([f"- {s}" for s in req.template_sections])
    
    # Build section-specific generation instructions
    section_instructions = {}
    section_instructions["Executive Summary"] = """Write a concise 1-page executive summary covering: project purpose, strategic alignment, expected outcomes, and call-to-action for vendors. Max 250 words. Decision-maker friendly."""
    section_instructions["Background & Context"] = """Describe the organizational context, current challenges, existing systems landscape, and the strategic drivers behind this procurement. Reference the knowledge base content. Include organizational structure if available."""
    section_instructions["Definitions & Acronyms"] = """Create a table of key terms and acronyms used in this RFP:
| Term | Definition |
| --- | --- |
List at least 10 relevant terms based on the project domain and source content."""
    section_instructions["Project Objectives"] = """Structure into 4 sub-categories with bullet points:
### Strategic Objectives
- (alignment with organizational strategy)
### Technical Objectives
- (systems, integrations, architecture goals)
### Operational Objectives
- (efficiency, workflow, process improvements)
### Governance Objectives
- (compliance, oversight, reporting)"""
    section_instructions["Scope of Work"] = """Group into 3-4 workstreams with a table per workstream:
| Activity | Description | Output |
| --- | --- | --- |
Include in-scope and out-of-scope statements."""
    section_instructions["Technical Requirements"] = """Organize into categories with bullet requirements:
### Integration Requirements
### Data Requirements
### Security Requirements
### Reporting & Analytics
3-4 bullets per category. Keep concise."""
    section_instructions["Deliverables & Acceptance Criteria"] = """Table format:
| Deliverable | Description | Format | Acceptance Criteria |
| --- | --- | --- | --- |
Include 6-8 deliverables covering design, dev, testing, training, docs, go-live."""
    section_instructions["Timeline & Milestones"] = """Define a phased timeline:
### Phase 1: Initiation & Planning (Weeks 1-4)
- Key milestones and governance gates
### Phase 2: Design & Development (Weeks 5-16)
- Key milestones
### Phase 3: Testing & UAT (Weeks 17-22)
- Key milestones
### Phase 4: Go-Live & Transition (Weeks 23-26)
- Key milestones
Include dependencies between phases and governance approval gates."""
    section_instructions["Commercial Model & Budget"] = """Define:
### Pricing Model
- Specify Fixed Price / Time & Materials / Hybrid expectations
### Cost Breakdown Structure
| Cost Category | Description | Pricing Type |
| --- | --- | --- |
Include categories: Professional Services, Software Licensing, Infrastructure, Training, Support & Maintenance, Project Management
### Payment Schedule
- Milestone-based payment structure tied to deliverables"""
    section_instructions["Evaluation Criteria"] = """Create a weighted scoring table:
| Criteria | Weight (%) | Sub-Criteria | Evaluation Method |
| --- | --- | --- | --- |
Include: Technical Approach (30-40%), Team & Experience (15-20%), Commercial (20-25%), Project Management (10-15%), Innovation (5-10%).
State the minimum technical threshold score for qualification."""
    section_instructions["Instructions to Bidders"] = """Include structured guidance:
### Communication Protocol
- Single point of contact, email format, response timelines
### Clarification Process
- Deadline for questions, method of submission, response distribution
### Proposal Validity
- Minimum validity period (e.g., 90 days)
### Conflict of Interest
- Declaration requirements and disqualification criteria"""
    section_instructions["Submission Requirements"] = """Specify:
### Proposal Structure
- Technical Proposal (separate volume)
- Financial Proposal (separate sealed volume)
### Naming Conventions
- File naming format requirements
### Format Requirements
- Page limits, font, margins, file types
### Deadline
- Exact submission date, time, timezone, and method"""
    section_instructions["Legal Terms & Conditions"] = """Cover:
### Liability & Indemnification
- Limitation of liability, indemnification obligations
### Confidentiality
- NDA requirements, information handling
### Intellectual Property
- IP ownership of deliverables, pre-existing IP, licensing
### Regulatory Compliance
- Applicable laws, regulations, standards
### Termination
- Termination for convenience/cause, notice periods"""
    section_instructions["Data Protection & Security Compliance"] = """Include:
### Healthcare Data Standards
- Applicable data protection regulations, patient data handling
### Access Control
- Role-based access, multi-factor authentication, least privilege
### Audit & Logging
- Audit trail requirements, log retention, monitoring
### Data Residency
- Data sovereignty requirements, hosting location
### Incident Response
- Breach notification timelines, response procedures"""
    section_instructions["Appendices"] = """Reference the following appendix items:
### Appendix A: System Integration List
- List all systems requiring integration with brief descriptions
### Appendix B: Data Model Templates
- Reference to data schemas and exchange formats
### Appendix C: Governance Framework
- Oversight structure, steering committee, escalation paths
### Appendix D: Risk Register Template
| Risk | Likelihood | Impact | Mitigation | Owner |
| --- | --- | --- | --- | --- |
Provide 5-8 sample risks relevant to the project."""

    # Build section-specific instructions for the batch
    batch_instructions = []
    for sec in req.template_sections:
        instr = section_instructions.get(sec, f"Write 2-3 substantive paragraphs with bullet points where appropriate. Structure content with sub-headings.")
        batch_instructions.append(f"=== {sec} ===\n{instr}")
    
    instructions_text = "\n\n".join(batch_instructions)
    
    prompt = f"""You are an Enterprise RFP Architect. Generate procurement-ready content for these sections:

{sections_list}

PROJECT: {req.project_name or 'Project'}
CLIENT: {req.client_name or 'Organization'}
TONE: {tone_map.get(req.tone, 'formal procurement')}
{f'CONTEXT: {req.additional_context}' if req.additional_context else ''}

KNOWLEDGE BASE:
{combined}

SECTION INSTRUCTIONS:
{instructions_text}

RULES:
- Use markdown tables (| Col | Col |), ### sub-headings, - bullets, **bold** where instructed
- No empty sections. No [Source:] citations. No paragraphs over 150 words without structure.
- Ground content in source material. Use project/client names naturally.
- If healthcare detected: include HIS, data governance, regulatory compliance.

Return ONLY JSON array: [{{"section":"Title","content":"..."}}]"""

    response = await generate_with_ai(prompt, f"Enterprise RFP writer. {req.tone} tone.")
    
    try:
        clean = response.strip()
        if clean.startswith('```'): clean = clean.split('\n', 1)[1].rsplit('```', 1)[0]
        rfp_sections = json.loads(clean)
    except:
        rfp_sections = [{"section": s, "content": f"Content for {s}."} for s in req.template_sections]
    
    return {"sections": rfp_sections, "source_names": source_names}


class RFPExportRequest(BaseModel):
    sections: List[dict]
    project_name: str = ""
    client_name: str = ""
    tone: str = "Formal"
    source_names: List[str] = []
    format: str = "docx"  # "docx" or "pdf"


@api_router.post("/rfp/export")
async def export_rfp(req: RFPExportRequest):
    """Export RFP as a professionally formatted DOCX or PDF"""
    import re
    date_str = datetime.now().strftime('%B %d, %Y')
    
    if req.format == "docx":
        return _export_docx(req, date_str)
    elif req.format == "pdf":
        return _export_pdf(req, date_str)
    else:
        raise HTTPException(status_code=400, detail="Format must be 'docx' or 'pdf'")


def _strip_citations(text: str) -> str:
    """Remove all [Source: ...] citation tags from text"""
    import re
    text = re.sub(r'\s*\[Source:[^\]]*\]\s*', ' ', text)
    text = re.sub(r'\s{2,}', ' ', text)
    return text.strip()


def _export_docx(req: RFPExportRequest, date_str: str):
    """Generate a professionally formatted Word document"""
    from docx import Document
    from docx.shared import Inches, Pt, Cm, RGBColor, Emu
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT
    from docx.oxml.ns import qn, nsdecls
    from docx.oxml import parse_xml
    import re
    
    doc = Document()
    
    # Page setup
    for section in doc.sections:
        section.top_margin = Cm(2.5)
        section.bottom_margin = Cm(2)
        section.left_margin = Cm(2.5)
        section.right_margin = Cm(2.5)
    
    # -- Define custom styles --
    # Normal style
    style = doc.styles['Normal']
    style.font.name = 'Calibri'
    style.font.size = Pt(11)
    style.font.color.rgb = RGBColor(55, 65, 81)
    style.paragraph_format.space_after = Pt(6)
    style.paragraph_format.line_spacing = 1.3
    
    # Heading 1 style (section headers)
    h1_style = doc.styles['Heading 1']
    h1_style.font.name = 'Calibri'
    h1_style.font.size = Pt(18)
    h1_style.font.bold = True
    h1_style.font.color.rgb = RGBColor(0, 77, 64)
    h1_style.paragraph_format.space_before = Pt(24)
    h1_style.paragraph_format.space_after = Pt(10)
    h1_style.paragraph_format.keep_with_next = True
    
    # Heading 2 style (TOC title)
    h2_style = doc.styles['Heading 2']
    h2_style.font.name = 'Calibri'
    h2_style.font.size = Pt(14)
    h2_style.font.bold = True
    h2_style.font.color.rgb = RGBColor(0, 60, 50)
    h2_style.paragraph_format.space_before = Pt(16)
    h2_style.paragraph_format.space_after = Pt(8)
    
    # ═══════════ TITLE PAGE ═══════════
    # Top spacing
    for _ in range(4):
        p = doc.add_paragraph()
        p.space_after = Pt(0)
    
    # Decorative top bar (using a table as a color block)
    tbl = doc.add_table(rows=1, cols=1)
    tbl.alignment = WD_TABLE_ALIGNMENT.CENTER
    cell = tbl.cell(0, 0)
    cell.width = Cm(15)
    shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="004D40"/>')
    cell.paragraphs[0]._element.get_or_add_pPr().append(shading)
    cell_para = cell.paragraphs[0]
    cell_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    cell_para.space_before = Pt(20)
    cell_para.space_after = Pt(20)
    run = cell_para.add_run('REQUEST FOR PROPOSAL')
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = RGBColor(200, 168, 107)
    run.font.name = 'Calibri'
    run.font.letter_spacing = Pt(3)
    
    doc.add_paragraph().space_after = Pt(16)
    
    # Project name
    name_para = doc.add_paragraph()
    name_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    name_para.space_after = Pt(6)
    run = name_para.add_run(req.project_name or 'Project')
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 77, 64)
    run.font.name = 'Calibri'
    
    # Gold divider
    div_para = doc.add_paragraph()
    div_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    div_para.space_before = Pt(8)
    div_para.space_after = Pt(8)
    run = div_para.add_run('\u2500' * 30)
    run.font.size = Pt(10)
    run.font.color.rgb = RGBColor(200, 168, 107)
    
    # Prepared for
    meta_para = doc.add_paragraph()
    meta_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    meta_para.space_after = Pt(2)
    run = meta_para.add_run('Prepared for')
    run.font.size = Pt(10)
    run.font.color.rgb = RGBColor(148, 163, 184)
    run.font.name = 'Calibri'
    
    client_para = doc.add_paragraph()
    client_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    client_para.space_after = Pt(16)
    run = client_para.add_run(req.client_name or 'Organization')
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = RGBColor(30, 30, 30)
    run.font.name = 'Calibri'
    
    # Details table (centered info block)
    doc.add_paragraph().space_after = Pt(4)
    info_tbl = doc.add_table(rows=3, cols=2)
    info_tbl.alignment = WD_TABLE_ALIGNMENT.CENTER
    info_data = [
        ('Date:', date_str),
        ('Document Type:', f'{req.tone} RFP'),
        ('Sections:', f'{len(req.sections)} sections'),
    ]
    for row_idx, (label, value) in enumerate(info_data):
        info_tbl.cell(row_idx, 0).width = Cm(4)
        info_tbl.cell(row_idx, 1).width = Cm(8)
        label_para = info_tbl.cell(row_idx, 0).paragraphs[0]
        label_para.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        run = label_para.add_run(label)
        run.font.size = Pt(10)
        run.font.color.rgb = RGBColor(120, 120, 120)
        run.font.bold = True
        val_para = info_tbl.cell(row_idx, 1).paragraphs[0]
        run = val_para.add_run(f'  {value}')
        run.font.size = Pt(10)
        run.font.color.rgb = RGBColor(55, 65, 81)
    # Remove table borders
    for row in info_tbl.rows:
        for cell in row.cells:
            tc = cell._element
            tcPr = tc.get_or_add_tcPr()
            tcBorders = parse_xml(f'<w:tcBorders {nsdecls("w")}><w:top w:val="none"/><w:left w:val="none"/><w:bottom w:val="none"/><w:right w:val="none"/></w:tcBorders>')
            tcPr.append(tcBorders)
    
    # Bottom bar
    for _ in range(4):
        p = doc.add_paragraph()
        p.space_after = Pt(0)
    
    bot_tbl = doc.add_table(rows=1, cols=1)
    bot_tbl.alignment = WD_TABLE_ALIGNMENT.CENTER
    bot_cell = bot_tbl.cell(0, 0)
    shading2 = parse_xml(f'<w:shd {nsdecls("w")} w:fill="004D40"/>')
    bot_cell.paragraphs[0]._element.get_or_add_pPr().append(shading2)
    bot_para = bot_cell.paragraphs[0]
    bot_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    bot_para.space_before = Pt(8)
    bot_para.space_after = Pt(8)
    run = bot_para.add_run('CONFIDENTIAL')
    run.font.size = Pt(9)
    run.font.color.rgb = RGBColor(200, 168, 107)
    run.font.bold = True
    run.font.letter_spacing = Pt(2)
    
    # ═══════════ TABLE OF CONTENTS ═══════════
    doc.add_page_break()
    
    toc_title = doc.add_heading('Table of Contents', level=2)
    for run in toc_title.runs:
        run.font.color.rgb = RGBColor(0, 77, 64)
        run.font.size = Pt(20)
    
    # Divider below TOC title
    toc_div = doc.add_paragraph()
    toc_div.space_after = Pt(12)
    run = toc_div.add_run('\u2500' * 50)
    run.font.size = Pt(8)
    run.font.color.rgb = RGBColor(200, 168, 107)
    
    # TOC entries as a table (number | title | dots)
    toc_tbl = doc.add_table(rows=len(req.sections), cols=2)
    toc_tbl.alignment = WD_TABLE_ALIGNMENT.LEFT
    for i, sec in enumerate(req.sections):
        # Number cell
        num_cell = toc_tbl.cell(i, 0)
        num_cell.width = Cm(1.5)
        num_para = num_cell.paragraphs[0]
        num_para.space_before = Pt(4)
        num_para.space_after = Pt(4)
        run = num_para.add_run(f'{i + 1}.')
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0, 77, 64)
        # Title cell
        title_cell = toc_tbl.cell(i, 1)
        title_cell.width = Cm(12)
        title_para = title_cell.paragraphs[0]
        title_para.space_before = Pt(4)
        title_para.space_after = Pt(4)
        run = title_para.add_run(sec.get("section", "Section"))
        run.font.size = Pt(12)
        run.font.color.rgb = RGBColor(55, 65, 81)
    # Remove table borders
    for row in toc_tbl.rows:
        for cell in row.cells:
            tc = cell._element
            tcPr = tc.get_or_add_tcPr()
            tcBorders = parse_xml(f'<w:tcBorders {nsdecls("w")}><w:top w:val="none"/><w:left w:val="none"/><w:bottom w:val="none"/><w:right w:val="none"/></w:tcBorders>')
            tcPr.append(tcBorders)
    # Add row bottom borders only
    for row in toc_tbl.rows:
        for cell in row.cells:
            tc = cell._element
            tcPr = tc.get_or_add_tcPr()
            bottom_border = parse_xml(f'<w:tcBorders {nsdecls("w")}><w:bottom w:val="single" w:sz="4" w:color="E5E7EB"/></w:tcBorders>')
            tcPr.append(bottom_border)
    
    # ═══════════ SECTIONS ═══════════
    for i, sec in enumerate(req.sections):
        doc.add_page_break()
        
        # Section heading using Heading 1 style (will appear in Word's built-in TOC)
        heading = doc.add_heading(f'{sec.get("section", "Section")}', level=1)
        heading.space_before = Pt(0)
        for run in heading.runs:
            run.font.color.rgb = RGBColor(0, 77, 64)
            run.font.size = Pt(18)
            run.font.name = 'Calibri'
        
        # Section number bar
        num_bar = doc.add_paragraph()
        num_bar.space_after = Pt(12)
        run = num_bar.add_run(f'Section {i + 1} of {len(req.sections)}')
        run.font.size = Pt(9)
        run.font.color.rgb = RGBColor(148, 163, 184)
        run.font.italic = True
        
        # Accent line under heading
        accent_line = doc.add_paragraph()
        accent_line.space_after = Pt(12)
        run = accent_line.add_run('\u2500' * 50)
        run.font.size = Pt(6)
        run.font.color.rgb = RGBColor(0, 77, 64)
        
        # Section content — parse structured markdown (tables, bullets, sub-headings)
        content = _strip_citations(sec.get('content', ''))
        lines = content.split('\n')
        li = 0
        while li < len(lines):
            line = lines[li].rstrip()
            
            # Markdown table
            if line.startswith('|') and '|' in line[1:]:
                table_lines = []
                while li < len(lines) and lines[li].strip().startswith('|'):
                    table_lines.append(lines[li].strip())
                    li += 1
                # Parse rows, skip separator lines (---|---)
                rows = []
                for tl in table_lines:
                    cells = [c.strip() for c in tl.split('|') if c.strip()]
                    if cells and not all(set(c) <= set('-: ') for c in cells):
                        rows.append(cells)
                if rows:
                    col_count = max(len(r) for r in rows)
                    tbl = doc.add_table(rows=len(rows), cols=col_count)
                    tbl.alignment = WD_TABLE_ALIGNMENT.LEFT
                    for ri, row_data in enumerate(rows):
                        for ci in range(col_count):
                            cell = tbl.cell(ri, ci)
                            val = row_data[ci] if ci < len(row_data) else ''
                            cell.paragraphs[0].text = ''
                            run = cell.paragraphs[0].add_run(val)
                            run.font.size = Pt(10)
                            run.font.name = 'Calibri'
                            if ri == 0:
                                run.font.bold = True
                                run.font.color.rgb = RGBColor(255, 255, 255)
                                shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="004D40"/>')
                                cell._element.get_or_add_tcPr().append(shading)
                            else:
                                run.font.color.rgb = RGBColor(55, 65, 81)
                                if ri % 2 == 0:
                                    shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="F8FAFB"/>')
                                    cell._element.get_or_add_tcPr().append(shading)
                    doc.add_paragraph().space_after = Pt(4)
                continue
            
            # Sub-heading (### or ##)
            if re.match(r'^#{2,4}\s', line):
                txt = re.sub(r'^#+\s*', '', line)
                h = doc.add_heading(txt, level=2)
                for run in h.runs:
                    run.font.color.rgb = RGBColor(0, 60, 50)
                    run.font.size = Pt(13)
                    run.font.name = 'Calibri'
                li += 1
                continue
            
            # Bullet point
            if re.match(r'^\s*[-*•]\s', line):
                bullets = []
                while li < len(lines) and re.match(r'^\s*[-*•]\s', lines[li]):
                    bullets.append(re.sub(r'^\s*[-*•]\s*', '', lines[li]))
                    li += 1
                for b in bullets:
                    p = doc.add_paragraph(style='List Bullet')
                    p.paragraph_format.left_indent = Cm(1)
                    p.space_after = Pt(3)
                    # Handle bold within bullet
                    parts = re.split(r'(\*\*[^*]+\*\*)', b)
                    for part in parts:
                        if part.startswith('**') and part.endswith('**'):
                            run = p.add_run(part[2:-2])
                            run.font.bold = True
                        else:
                            run = p.add_run(part)
                        run.font.size = Pt(11)
                        run.font.color.rgb = RGBColor(55, 65, 81)
                        run.font.name = 'Calibri'
                continue
            
            # Numbered list
            if re.match(r'^\s*\d+[.)]\s', line):
                items = []
                while li < len(lines) and re.match(r'^\s*\d+[.)]\s', lines[li]):
                    items.append(re.sub(r'^\s*\d+[.)]\s*', '', lines[li]))
                    li += 1
                for item in items:
                    p = doc.add_paragraph(style='List Number')
                    p.paragraph_format.left_indent = Cm(1)
                    p.space_after = Pt(3)
                    parts = re.split(r'(\*\*[^*]+\*\*)', item)
                    for part in parts:
                        if part.startswith('**') and part.endswith('**'):
                            run = p.add_run(part[2:-2])
                            run.font.bold = True
                        else:
                            run = p.add_run(part)
                        run.font.size = Pt(11)
                        run.font.color.rgb = RGBColor(55, 65, 81)
                        run.font.name = 'Calibri'
                continue
            
            # Empty line
            if not line.strip():
                li += 1
                continue
            
            # Regular paragraph with bold support
            para = doc.add_paragraph()
            para.space_after = Pt(8)
            para.paragraph_format.line_spacing = 1.3
            parts = re.split(r'(\*\*[^*]+\*\*)', line)
            for part in parts:
                if part.startswith('**') and part.endswith('**'):
                    run = para.add_run(part[2:-2])
                    run.font.bold = True
                else:
                    run = para.add_run(part)
                run.font.size = Pt(11)
                run.font.color.rgb = RGBColor(55, 65, 81)
                run.font.name = 'Calibri'
            li += 1
    
    # ═══════════ FOOTER PAGE ═══════════
    doc.add_page_break()
    for _ in range(6):
        p = doc.add_paragraph()
        p.space_after = Pt(0)
    
    end_para = doc.add_paragraph()
    end_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = end_para.add_run('End of Document')
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(148, 163, 184)
    run.font.italic = True
    
    end_div = doc.add_paragraph()
    end_div.alignment = WD_ALIGN_PARAGRAPH.CENTER
    end_div.space_before = Pt(8)
    run = end_div.add_run('\u2500' * 30)
    run.font.size = Pt(8)
    run.font.color.rgb = RGBColor(200, 168, 107)
    
    gen_para = doc.add_paragraph()
    gen_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    gen_para.space_before = Pt(8)
    run = gen_para.add_run('Generated by Knowledge Base  |  AI-Powered Research Assistant')
    run.font.size = Pt(9)
    run.font.color.rgb = RGBColor(148, 163, 184)
    
    # Save to buffer
    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    
    import re as _re
    safe_name = _re.sub(r'[^\x20-\x7E]', '', (req.project_name or 'RFP')).replace(' ', '_') or 'RFP'
    filename = f"{safe_name}_RFP.docx"
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'}
    )


def _export_pdf(req: RFPExportRequest, date_str: str):
    """Generate a professionally formatted PDF document"""
    from fpdf import FPDF
    import re
    
    TEAL = (0, 77, 64)
    GOLD = (200, 168, 107)
    DARK = (30, 31, 54)
    GRAY = (55, 65, 81)
    LIGHT = (148, 163, 184)
    WHITE = (255, 255, 255)
    
    class RFPPDF(FPDF):
        def header(self):
            if self.page_no() > 2:  # Skip header on title + TOC pages
                self.set_font('Helvetica', 'B', 8)
                self.set_text_color(*TEAL)
                self.cell(0, 6, self.project_name, align='L')
                self.set_font('Helvetica', '', 8)
                self.set_text_color(*LIGHT)
                self.cell(0, 6, f'Page {self.page_no()}', align='R', new_x="LMARGIN", new_y="NEXT")
                self.set_draw_color(*GOLD)
                self.set_line_width(0.4)
                self.line(self.l_margin, self.get_y() + 1, self.w - self.r_margin, self.get_y() + 1)
                self.ln(6)
        
        def footer(self):
            self.set_y(-12)
            self.set_draw_color(*GOLD)
            self.set_line_width(0.3)
            self.line(self.l_margin, self.get_y(), self.w - self.r_margin, self.get_y())
            self.ln(2)
            self.set_font('Helvetica', '', 7)
            self.set_text_color(*LIGHT)
            self.cell(0, 6, f'{self.project_name}  |  Confidential  |  {self.date_str}', align='C')
    
    pdf = RFPPDF()
    pdf.project_name = req.project_name or 'Project'
    pdf.date_str = date_str
    pdf.set_auto_page_break(auto=True, margin=18)
    pdf.set_margins(25, 20, 25)
    pw = pdf.w - 50  # printable width
    
    # ═══════════ TITLE PAGE ═══════════
    pdf.add_page()
    
    # Top teal banner
    pdf.set_fill_color(*TEAL)
    pdf.rect(0, 0, pdf.w, 10, 'F')
    # Gold accent line below banner
    pdf.set_fill_color(*GOLD)
    pdf.rect(0, 10, pdf.w, 2, 'F')
    
    pdf.ln(55)
    
    # "REQUEST FOR PROPOSAL" label
    pdf.set_font('Helvetica', 'B', 11)
    pdf.set_text_color(*GOLD)
    pdf.cell(0, 7, 'REQUEST FOR PROPOSAL', align='C', new_x="LMARGIN", new_y="NEXT")
    pdf.ln(8)
    
    # Project name
    pdf.set_font('Helvetica', 'B', 30)
    pdf.set_text_color(*TEAL)
    pdf.multi_cell(0, 13, req.project_name or 'Project', align='C')
    pdf.ln(6)
    
    # Gold divider
    cx = pdf.w / 2
    pdf.set_draw_color(*GOLD)
    pdf.set_line_width(1)
    pdf.line(cx - 35, pdf.get_y(), cx + 35, pdf.get_y())
    pdf.ln(12)
    
    # "Prepared for"
    pdf.set_font('Helvetica', '', 10)
    pdf.set_text_color(*LIGHT)
    pdf.cell(0, 6, 'Prepared for', align='C', new_x="LMARGIN", new_y="NEXT")
    pdf.ln(2)
    pdf.set_font('Helvetica', 'B', 18)
    pdf.set_text_color(*DARK)
    pdf.cell(0, 10, req.client_name or 'Organization', align='C', new_x="LMARGIN", new_y="NEXT")
    pdf.ln(16)
    
    # Details box
    box_x = pdf.w / 2 - 45
    box_y = pdf.get_y()
    pdf.set_draw_color(220, 220, 220)
    pdf.set_line_width(0.3)
    pdf.rect(box_x, box_y, 90, 40, 'D')
    # Teal top edge on box
    pdf.set_fill_color(*TEAL)
    pdf.rect(box_x, box_y, 90, 2, 'F')
    
    pdf.set_y(box_y + 6)
    pdf.set_font('Helvetica', '', 9)
    pdf.set_text_color(*LIGHT)
    pdf.cell(0, 5, 'Date', align='C', new_x="LMARGIN", new_y="NEXT")
    pdf.set_font('Helvetica', 'B', 10)
    pdf.set_text_color(*GRAY)
    pdf.cell(0, 6, date_str, align='C', new_x="LMARGIN", new_y="NEXT")
    pdf.ln(2)
    pdf.set_font('Helvetica', '', 9)
    pdf.set_text_color(*LIGHT)
    pdf.cell(0, 5, f'Document Type: {req.tone} RFP  |  {len(req.sections)} Sections', align='C', new_x="LMARGIN", new_y="NEXT")
    
    # Bottom teal banner
    pdf.set_fill_color(*TEAL)
    pdf.rect(0, pdf.h - 12, pdf.w, 10, 'F')
    pdf.set_fill_color(*GOLD)
    pdf.rect(0, pdf.h - 14, pdf.w, 2, 'F')
    pdf.set_y(pdf.h - 11)
    pdf.set_font('Helvetica', 'B', 7)
    pdf.set_text_color(*GOLD)
    pdf.cell(0, 8, 'CONFIDENTIAL', align='C')
    
    # ═══════════ TABLE OF CONTENTS ═══════════
    pdf.add_page()
    
    # TOC header
    pdf.set_font('Helvetica', 'B', 22)
    pdf.set_text_color(*TEAL)
    pdf.cell(0, 14, 'Table of Contents', new_x="LMARGIN", new_y="NEXT")
    
    # Gold underline
    pdf.set_draw_color(*GOLD)
    pdf.set_line_width(1)
    pdf.line(pdf.l_margin, pdf.get_y() + 2, pdf.l_margin + 60, pdf.get_y() + 2)
    pdf.ln(10)
    
    # TOC entries with alternating background
    for i, sec in enumerate(req.sections):
        y = pdf.get_y()
        if i % 2 == 0:
            pdf.set_fill_color(244, 246, 251)
            pdf.rect(pdf.l_margin, y, pw, 10, 'F')
        
        # Section number (teal bold)
        pdf.set_font('Helvetica', 'B', 11)
        pdf.set_text_color(*TEAL)
        pdf.set_xy(pdf.l_margin + 4, y)
        pdf.cell(12, 10, f'{i + 1}.')
        
        # Section title
        pdf.set_font('Helvetica', '', 11)
        pdf.set_text_color(*GRAY)
        pdf.set_xy(pdf.l_margin + 18, y)
        pdf.cell(pw - 18, 10, sec.get("section", "Section"), new_x="LMARGIN", new_y="NEXT")
    
    # ═══════════ SECTIONS ═══════════
    for i, sec in enumerate(req.sections):
        pdf.add_page()
        
        # Section header with teal background
        header_y = pdf.get_y()
        pdf.set_fill_color(*TEAL)
        pdf.rect(pdf.l_margin, header_y, pw, 18, 'F')
        # Gold accent
        pdf.set_fill_color(*GOLD)
        pdf.rect(pdf.l_margin, header_y, 4, 18, 'F')
        
        # Section number + title (white on teal)
        pdf.set_font('Helvetica', 'B', 14)
        pdf.set_text_color(*WHITE)
        pdf.set_xy(pdf.l_margin + 10, header_y + 2)
        pdf.cell(pw - 10, 14, f'{i + 1}.  {sec.get("section", "Section")}')
        
        pdf.set_y(header_y + 22)
        
        # Section subtitle
        pdf.set_font('Helvetica', 'I', 8)
        pdf.set_text_color(*LIGHT)
        pdf.cell(0, 5, f'Section {i + 1} of {len(req.sections)}', new_x="LMARGIN", new_y="NEXT")
        pdf.ln(4)
        
        # Content — strip all citations
        content = _strip_citations(sec.get('content', ''))
        paragraphs = content.split('\n\n')
        
        for para_text in paragraphs:
            para_text = para_text.strip()
            if not para_text:
                continue
            pdf.set_font('Helvetica', '', 10.5)
            pdf.set_text_color(*GRAY)
            pdf.multi_cell(pw, 5.5, para_text, align='J')
            pdf.ln(4)
    
    # ═══════════ END PAGE ═══════════
    pdf.add_page()
    pdf.ln(60)
    pdf.set_font('Helvetica', 'I', 14)
    pdf.set_text_color(*LIGHT)
    pdf.cell(0, 10, 'End of Document', align='C', new_x="LMARGIN", new_y="NEXT")
    pdf.ln(4)
    pdf.set_draw_color(*GOLD)
    pdf.set_line_width(0.8)
    pdf.line(cx - 25, pdf.get_y(), cx + 25, pdf.get_y())
    pdf.ln(8)
    pdf.set_font('Helvetica', '', 9)
    pdf.cell(0, 6, 'Generated by Knowledge Base  |  AI-Powered Research Assistant', align='C')
    
    # Save to buffer
    buf = io.BytesIO()
    pdf.output(buf)
    buf.seek(0)
    
    import re as _re
    safe_name = _re.sub(r'[^\x20-\x7E]', '', (req.project_name or 'RFP')).replace(' ', '_') or 'RFP'
    filename = f"{safe_name}_RFP.pdf"
    return StreamingResponse(
        buf,
        media_type="application/pdf",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'}
    )

# ── Quiz DOCX Export ──────────────────────────────────────────────────────
class QuizExportRequest(BaseModel):
    title: str
    quiz_data: dict  # {mcq:[], truefalse:[], short:[]}

@api_router.post("/export/quiz-docx")
async def export_quiz_docx(req: QuizExportRequest):
    """Export quiz as a professionally formatted Word document"""
    from docx import Document
    from docx.shared import Pt, Cm, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT
    from docx.oxml.ns import nsdecls
    from docx.oxml import parse_xml

    doc = Document()
    for section in doc.sections:
        section.top_margin = Cm(2.5)
        section.bottom_margin = Cm(2)
        section.left_margin = Cm(2.5)
        section.right_margin = Cm(2.5)

    style = doc.styles['Normal']
    style.font.name = 'Calibri'
    style.font.size = Pt(11)
    style.font.color.rgb = RGBColor(55, 65, 81)
    style.paragraph_format.line_spacing = 1.3

    h1 = doc.styles['Heading 1']
    h1.font.name = 'Calibri'
    h1.font.size = Pt(18)
    h1.font.bold = True
    h1.font.color.rgb = RGBColor(0, 77, 64)
    h1.paragraph_format.space_before = Pt(24)
    h1.paragraph_format.space_after = Pt(10)

    h2 = doc.styles['Heading 2']
    h2.font.name = 'Calibri'
    h2.font.size = Pt(14)
    h2.font.bold = True
    h2.font.color.rgb = RGBColor(0, 60, 50)
    h2.paragraph_format.space_before = Pt(16)
    h2.paragraph_format.space_after = Pt(8)

    mcq = req.quiz_data.get('mcq', [])
    tf = req.quiz_data.get('truefalse', [])
    short = req.quiz_data.get('short', [])
    total = len(mcq) + len(tf) + len(short)

    # Title page
    for _ in range(3):
        doc.add_paragraph().space_after = Pt(0)

    tbl = doc.add_table(rows=1, cols=1)
    tbl.alignment = WD_TABLE_ALIGNMENT.CENTER
    cell = tbl.cell(0, 0)
    shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="004D40"/>')
    cell.paragraphs[0]._element.get_or_add_pPr().append(shading)
    cell_p = cell.paragraphs[0]
    cell_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    cell_p.space_before = Pt(20)
    cell_p.space_after = Pt(20)
    run = cell_p.add_run('KNOWLEDGE ASSESSMENT')
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = RGBColor(200, 168, 107)
    run.font.name = 'Calibri'

    doc.add_paragraph().space_after = Pt(12)

    tp = doc.add_paragraph()
    tp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = tp.add_run(req.title)
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 77, 64)

    div = doc.add_paragraph()
    div.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = div.add_run('\u2500' * 30)
    run.font.color.rgb = RGBColor(200, 168, 107)
    run.font.size = Pt(10)

    meta = doc.add_paragraph()
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    meta.space_after = Pt(4)
    run = meta.add_run(f'{total} Questions  |  MCQ: {len(mcq)}  |  True/False: {len(tf)}  |  Short Answer: {len(short)}')
    run.font.size = Pt(11)
    run.font.color.rgb = RGBColor(100, 116, 139)

    # Info fields
    doc.add_paragraph().space_after = Pt(8)
    for field in ['Name: ____________________________', 'Date: ____________', f'Score: ______ / {total}']:
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.space_after = Pt(2)
        run = p.add_run(field)
        run.font.size = Pt(11)
        run.font.color.rgb = RGBColor(100, 116, 139)

    doc.add_page_break()
    qnum = 0

    # MCQ Section
    if mcq:
        doc.add_heading('Multiple Choice Questions', level=1)
        for q in mcq:
            qnum += 1
            p = doc.add_paragraph()
            p.space_before = Pt(10)
            run = p.add_run(f'{qnum}. ')
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 77, 64)
            run = p.add_run(q.get('question', ''))
            run.font.bold = True
            for opt in q.get('options', []):
                op = doc.add_paragraph(style='List Bullet')
                op.paragraph_format.left_indent = Cm(1.2)
                op.space_after = Pt(2)
                run = op.add_run(opt)
                run.font.size = Pt(11)
            doc.add_paragraph().space_after = Pt(4)

    # True/False Section
    if tf:
        doc.add_heading('True or False', level=1)
        for q in tf:
            qnum += 1
            p = doc.add_paragraph()
            p.space_before = Pt(8)
            run = p.add_run(f'{qnum}. ')
            run.font.bold = True
            run.font.color.rgb = RGBColor(14, 165, 233)
            run = p.add_run(q.get('statement', ''))
            run.font.bold = True
            tp = doc.add_paragraph()
            tp.paragraph_format.left_indent = Cm(1.2)
            tp.space_after = Pt(4)
            run = tp.add_run('True  /  False       (circle one)')
            run.font.size = Pt(10)
            run.font.color.rgb = RGBColor(148, 163, 184)

    # Short Answer Section
    if short:
        doc.add_heading('Short Answer Questions', level=1)
        for q in short:
            qnum += 1
            p = doc.add_paragraph()
            p.space_before = Pt(10)
            run = p.add_run(f'{qnum}. ')
            run.font.bold = True
            run.font.color.rgb = RGBColor(139, 92, 246)
            run = p.add_run(q.get('question', ''))
            run.font.bold = True
            for _ in range(3):
                lp = doc.add_paragraph()
                lp.space_after = Pt(0)
                run = lp.add_run('_' * 80)
                run.font.size = Pt(9)
                run.font.color.rgb = RGBColor(229, 231, 235)

    # Answer Key
    doc.add_page_break()
    doc.add_heading('Answer Key', level=1)
    anum = 0
    if mcq:
        doc.add_heading('Multiple Choice', level=2)
        for q in mcq:
            anum += 1
            p = doc.add_paragraph()
            p.space_after = Pt(2)
            run = p.add_run(f'{anum}. {q.get("correct", "—")}')
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 77, 64)
            if q.get('explanation'):
                run = p.add_run(f'  —  {q["explanation"]}')
                run.font.color.rgb = RGBColor(100, 116, 139)
                run.font.size = Pt(10)
    if tf:
        doc.add_heading('True / False', level=2)
        for q in tf:
            anum += 1
            p = doc.add_paragraph()
            p.space_after = Pt(2)
            ans = 'True' if q.get('answer') else 'False'
            run = p.add_run(f'{anum}. {ans}')
            run.font.bold = True
            run.font.color.rgb = RGBColor(14, 165, 233)
            if q.get('explanation'):
                run = p.add_run(f'  —  {q["explanation"]}')
                run.font.color.rgb = RGBColor(100, 116, 139)
                run.font.size = Pt(10)
    if short:
        doc.add_heading('Short Answer', level=2)
        for q in short:
            anum += 1
            p = doc.add_paragraph()
            p.space_after = Pt(4)
            run = p.add_run(f'{anum}. ')
            run.font.bold = True
            run.font.color.rgb = RGBColor(139, 92, 246)
            run = p.add_run(q.get('sampleAnswer', ''))
            run.font.size = Pt(10)
            run.font.color.rgb = RGBColor(100, 116, 139)

    # Footer
    p = doc.add_paragraph()
    p.space_before = Pt(30)
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run('Generated by Knowledge Base  |  AI-Powered Research Assistant')
    run.font.size = Pt(9)
    run.font.color.rgb = RGBColor(148, 163, 184)

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    import re as _re
    safe_title = _re.sub(r'[^\x20-\x7E]', '', req.title).replace(' ', '_') or 'quiz'
    filename = f'{safe_title}_quiz.docx'
    return StreamingResponse(buf, media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                             headers={"Content-Disposition": f'attachment; filename="{filename}"'})


# ── Data Table Excel Export ──────────────────────────────────────────────
class DataTableExportRequest(BaseModel):
    title: str
    table_data: dict  # {tables:[], stats:[]}

@api_router.post("/export/datatable-xlsx")
async def export_datatable_xlsx(req: DataTableExportRequest):
    """Export data table as a professionally formatted Excel file"""
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    wb = openpyxl.Workbook()
    tables = req.table_data.get('tables', [])
    stats = req.table_data.get('stats', [])

    teal_fill = PatternFill(start_color='004D40', end_color='004D40', fill_type='solid')
    gold_font = Font(name='Calibri', size=11, bold=True, color='C8A86B')
    header_font = Font(name='Calibri', size=11, bold=True, color='FFFFFF')
    title_font = Font(name='Calibri', size=14, bold=True, color='004D40')
    stat_val_font = Font(name='Calibri', size=16, bold=True, color='004D40')
    stat_lbl_font = Font(name='Calibri', size=10, color='64748B')
    data_font = Font(name='Calibri', size=11, color='374151')
    thin_border = Border(
        left=Side(style='thin', color='E5E7EB'),
        right=Side(style='thin', color='E5E7EB'),
        top=Side(style='thin', color='E5E7EB'),
        bottom=Side(style='thin', color='E5E7EB')
    )

    # Summary sheet
    ws = wb.active
    ws.title = 'Summary'
    ws.merge_cells('A1:E1')
    ws['A1'] = req.title
    ws['A1'].font = Font(name='Calibri', size=18, bold=True, color='004D40')
    ws['A1'].alignment = Alignment(horizontal='center')
    ws['A2'] = f'{len(tables)} Tables  |  {len(stats)} Key Metrics  |  Knowledge Base'
    ws.merge_cells('A2:E2')
    ws['A2'].font = Font(name='Calibri', size=10, color='94A3B8')
    ws['A2'].alignment = Alignment(horizontal='center')

    if stats:
        row = 4
        ws.cell(row=row, column=1, value='KEY METRICS').font = Font(name='Calibri', size=12, bold=True, color='004D40')
        row += 1
        for si, stat in enumerate(stats):
            ws.cell(row=row, column=1, value=stat.get('label', '')).font = stat_lbl_font
            ws.cell(row=row, column=2, value=stat.get('value', '')).font = stat_val_font
            ws.cell(row=row, column=3, value=stat.get('description', '')).font = stat_lbl_font
            row += 1
        row += 1
        ws.cell(row=row, column=1, value='TABLES').font = Font(name='Calibri', size=12, bold=True, color='004D40')
        row += 1
        for ti, table in enumerate(tables):
            ws.cell(row=row, column=1, value=f'{ti + 1}. {table.get("title", "Table")}').font = data_font
            ws.cell(row=row, column=2, value=f'{len(table.get("rows", []))} rows').font = stat_lbl_font
            row += 1

    ws.column_dimensions['A'].width = 25
    ws.column_dimensions['B'].width = 20
    ws.column_dimensions['C'].width = 40

    # Table sheets
    for ti, table in enumerate(tables):
        sheet_name = (table.get('title', f'Table {ti+1}'))[:31]
        wst = wb.create_sheet(title=sheet_name)
        headers = table.get('headers', [])
        rows = table.get('rows', [])

        # Title
        if headers:
            wst.merge_cells(start_row=1, start_column=1, end_row=1, end_column=len(headers))
        wst.cell(row=1, column=1, value=table.get('title', f'Table {ti+1}')).font = title_font
        wst.cell(row=1, column=1).alignment = Alignment(horizontal='left')

        # Headers
        for ci, h in enumerate(headers):
            cell = wst.cell(row=3, column=ci + 1, value=h)
            cell.font = header_font
            cell.fill = teal_fill
            cell.alignment = Alignment(horizontal='center', vertical='center')
            cell.border = thin_border
            wst.column_dimensions[get_column_letter(ci + 1)].width = max(15, len(str(h)) + 5)

        # Data rows
        alt_fill = PatternFill(start_color='F8FAFB', end_color='F8FAFB', fill_type='solid')
        for ri, row_data in enumerate(rows):
            if not isinstance(row_data, list):
                continue
            for ci, val in enumerate(row_data):
                cell = wst.cell(row=4 + ri, column=ci + 1, value=str(val) if val is not None else '')
                cell.font = data_font
                cell.border = thin_border
                cell.alignment = Alignment(vertical='center')
                if ri % 2 == 0:
                    cell.fill = alt_fill

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    import re as _re
    safe_title = _re.sub(r'[^\x20-\x7E]', '', req.title).replace(' ', '_') or 'data'
    filename = f'{safe_title}_data.xlsx'
    return StreamingResponse(buf, media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                             headers={"Content-Disposition": f'attachment; filename="{filename}"'})


# Include the router in the main app
app.include_router(api_router)

app.add_middleware(
    CORSMiddleware,
    allow_credentials=True,
    allow_origins=os.environ.get('CORS_ORIGINS', '*').split(','),
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.on_event("shutdown")
async def shutdown_db_client():
    client.close()
